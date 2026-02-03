"""QA validator — inspects generated PPTX output against a TemplateSchema.

Validates that a built presentation matches its schema contract: correct
slide count, data slots populated, formatting rules applied, table row
counts matching data, and chart series rendered.  Uses python-pptx to
read back the generated file.

Usage::

    from src.qa.validator import QAValidator

    validator = QAValidator(schema)
    result = validator.validate(pptx_bytes, payload)
    assert result.passed, result.summary()
"""

import io
import math
import re
from dataclasses import dataclass, field
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from src.schema.design_system import format_value, variance_color
from src.schema.models import (
    ChartType,
    DataSlot,
    FormatType,
    SlideSchema,
    SlideType,
    SlotType,
    TemplateSchema,
)


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------

@dataclass
class Issue:
    """A single QA issue found during validation."""
    severity: str       # "error" or "warning"
    slide_index: int    # -1 for presentation-level issues
    slide_name: str
    slot_name: str      # "" for slide-level issues
    category: str       # e.g. "slide_count", "slot_missing", "format"
    message: str

    def __str__(self) -> str:
        loc = f"slide {self.slide_index}"
        if self.slide_name:
            loc += f" ({self.slide_name})"
        if self.slot_name:
            loc += f" / {self.slot_name}"
        return f"[{self.severity.upper()}] {loc}: {self.message}"


@dataclass
class QAResult:
    """Aggregated result of QA validation."""
    issues: list[Issue] = field(default_factory=list)

    @property
    def errors(self) -> list[Issue]:
        return [i for i in self.issues if i.severity == "error"]

    @property
    def warnings(self) -> list[Issue]:
        return [i for i in self.issues if i.severity == "warning"]

    @property
    def passed(self) -> bool:
        return len(self.errors) == 0

    @property
    def error_count(self) -> int:
        return len(self.errors)

    @property
    def warning_count(self) -> int:
        return len(self.warnings)

    def summary(self) -> str:
        """One-line summary string."""
        status = "PASS" if self.passed else "FAIL"
        return (
            f"QA {status}: {self.error_count} error(s), "
            f"{self.warning_count} warning(s)"
        )

    def report(self) -> str:
        """Multi-line report of all issues."""
        lines = [self.summary()]
        for issue in self.issues:
            lines.append(f"  {issue}")
        return "\n".join(lines)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _is_missing(value: Any) -> bool:
    """Check if a value is None or NaN."""
    if value is None:
        return True
    if isinstance(value, float) and math.isnan(value):
        return True
    return False


def _all_text_on_slide(slide) -> str:
    """Concatenate all text on a slide for content searches."""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return " ".join(parts)


def _table_shapes(slide) -> list:
    """Return all table shapes on a slide."""
    return [s for s in slide.shapes if s.has_table]


def _chart_shapes(slide) -> list:
    """Return all chart shapes on a slide."""
    return [s for s in slide.shapes if s.has_chart]


# ---------------------------------------------------------------------------
# QAValidator
# ---------------------------------------------------------------------------

class QAValidator:
    """Validates generated PPTX output against a TemplateSchema.

    Parameters
    ----------
    schema : TemplateSchema
        The schema that was used to generate the presentation.
    """

    def __init__(self, schema: TemplateSchema) -> None:
        self.schema = schema

    def validate(self, pptx_bytes: bytes,
                 payload: dict[str, Any]) -> QAResult:
        """Run all validation checks on a built PPTX.

        Parameters
        ----------
        pptx_bytes : bytes
            The raw PPTX file content (from PPTXBuilder.build()).
        payload : dict[str, Any]
            The data payload that was used to build the presentation.

        Returns
        -------
        QAResult
            Aggregated validation result.
        """
        prs = Presentation(io.BytesIO(pptx_bytes))
        result = QAResult()

        self._check_slide_count(prs, result)
        self._check_dimensions(prs, result)
        self._check_payload_coverage(payload, result)

        # Per-slide checks (only if count matches)
        if len(prs.slides) == len(self.schema.slides):
            for slide_schema in self.schema.slides:
                slide = prs.slides[slide_schema.index]
                self._check_slide(slide, slide_schema, payload, result)

        return result

    def validate_payload(self, payload: dict[str, Any]) -> QAResult:
        """Validate a data payload against the schema without a PPTX.

        Checks that all required data keys are present, data types
        are correct, and table/chart data is structurally valid.
        """
        result = QAResult()
        self._check_payload_coverage(payload, result)
        self._check_payload_types(payload, result)
        return result

    # ------------------------------------------------------------------
    # Presentation-level checks
    # ------------------------------------------------------------------

    def _check_slide_count(self, prs: Presentation, result: QAResult) -> None:
        """Verify slide count matches schema."""
        expected = len(self.schema.slides)
        actual = len(prs.slides)
        if actual != expected:
            result.issues.append(Issue(
                severity="error",
                slide_index=-1,
                slide_name="",
                slot_name="",
                category="slide_count",
                message=f"Expected {expected} slides, got {actual}",
            ))

    def _check_dimensions(self, prs: Presentation, result: QAResult) -> None:
        """Verify presentation dimensions match schema."""
        expected_w = Inches(self.schema.width_inches)
        expected_h = Inches(self.schema.height_inches)
        if prs.slide_width != expected_w:
            result.issues.append(Issue(
                severity="error",
                slide_index=-1,
                slide_name="",
                slot_name="",
                category="dimensions",
                message=(
                    f"Slide width {prs.slide_width} != "
                    f"expected {expected_w}"
                ),
            ))
        if prs.slide_height != expected_h:
            result.issues.append(Issue(
                severity="error",
                slide_index=-1,
                slide_name="",
                slot_name="",
                category="dimensions",
                message=(
                    f"Slide height {prs.slide_height} != "
                    f"expected {expected_h}"
                ),
            ))

    # ------------------------------------------------------------------
    # Payload coverage checks
    # ------------------------------------------------------------------

    def _check_payload_coverage(self, payload: dict[str, Any],
                                result: QAResult) -> None:
        """Check that all schema data keys are present in the payload."""
        required = self.schema.all_data_keys()
        # Exclude column-level data_keys (those are within row dicts)
        top_level_keys: set[str] = set()
        for slide in self.schema.slides:
            for slot in slide.slots:
                top_level_keys.add(slot.data_key)
                if slot.variance_key:
                    top_level_keys.add(slot.variance_key)
                if slot.row_data_key:
                    top_level_keys.add(slot.row_data_key)
                if slot.categories_key:
                    top_level_keys.add(slot.categories_key)
                for series in slot.series:
                    top_level_keys.add(series.data_key)

        missing = top_level_keys - set(payload.keys())
        for key in sorted(missing):
            # Find which slide this key belongs to
            slide_name = self._find_slide_for_key(key)
            result.issues.append(Issue(
                severity="warning",
                slide_index=-1,
                slide_name=slide_name,
                slot_name="",
                category="payload_missing",
                message=f"Data key '{key}' not in payload",
            ))

    def _check_payload_types(self, payload: dict[str, Any],
                             result: QAResult) -> None:
        """Validate data types in the payload match schema expectations."""
        for slide_schema in self.schema.slides:
            for slot in slide_schema.slots:
                self._check_slot_payload(
                    slot, slide_schema, payload, result,
                )

    def _check_slot_payload(self, slot: DataSlot, slide_schema: SlideSchema,
                            payload: dict[str, Any],
                            result: QAResult) -> None:
        """Check a single slot's data in the payload."""
        value = payload.get(slot.data_key)

        # Table slots: row_data_key should be a list of dicts
        if slot.slot_type == SlotType.TABLE and slot.row_data_key:
            rows = payload.get(slot.row_data_key)
            if rows is not None and not isinstance(rows, list):
                result.issues.append(Issue(
                    severity="error",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="type_error",
                    message=(
                        f"row_data_key '{slot.row_data_key}' should be "
                        f"a list, got {type(rows).__name__}"
                    ),
                ))
            elif isinstance(rows, list) and rows:
                # Check column keys are present in row dicts
                row_keys = set(rows[0].keys()) if rows else set()
                for col in slot.columns:
                    if col.data_key not in row_keys:
                        result.issues.append(Issue(
                            severity="warning",
                            slide_index=slide_schema.index,
                            slide_name=slide_schema.name,
                            slot_name=slot.name,
                            category="column_key_missing",
                            message=(
                                f"Column '{col.header}' expects key "
                                f"'{col.data_key}' not found in row data"
                            ),
                        ))

        # Chart slots: series data should be lists
        if slot.slot_type == SlotType.CHART and slot.series:
            categories = payload.get(slot.categories_key) if slot.categories_key else None
            for series in slot.series:
                series_data = payload.get(series.data_key)
                if series_data is None:
                    continue
                if slot.chart_type in (ChartType.DOUGHNUT, ChartType.DOUGHNUT_EXPLODED):
                    # Doughnut series values are scalars
                    continue
                if not isinstance(series_data, (list, tuple)):
                    result.issues.append(Issue(
                        severity="error",
                        slide_index=slide_schema.index,
                        slide_name=slide_schema.name,
                        slot_name=slot.name,
                        category="type_error",
                        message=(
                            f"Series '{series.name}' data_key "
                            f"'{series.data_key}' should be a list, "
                            f"got {type(series_data).__name__}"
                        ),
                    ))
                elif categories and isinstance(categories, (list, tuple)):
                    if len(series_data) != len(categories):
                        result.issues.append(Issue(
                            severity="error",
                            slide_index=slide_schema.index,
                            slide_name=slide_schema.name,
                            slot_name=slot.name,
                            category="series_length_mismatch",
                            message=(
                                f"Series '{series.name}' has "
                                f"{len(series_data)} values but "
                                f"{len(categories)} categories"
                            ),
                        ))

        # KPI slots: value should be numeric (or None)
        if slot.slot_type == SlotType.KPI_VALUE and value is not None:
            if not isinstance(value, (int, float, str)):
                result.issues.append(Issue(
                    severity="error",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="type_error",
                    message=(
                        f"KPI value for '{slot.data_key}' should be "
                        f"numeric or string, got {type(value).__name__}"
                    ),
                ))

    # ------------------------------------------------------------------
    # Per-slide checks
    # ------------------------------------------------------------------

    def _check_slide(self, slide, slide_schema: SlideSchema,
                     payload: dict[str, Any], result: QAResult) -> None:
        """Run all checks for a single slide."""
        if slide_schema.slide_type == SlideType.SECTION_DIVIDER:
            self._check_divider_background(
                slide, slide_schema, result,
            )

        for slot in slide_schema.slots:
            self._check_slot(slide, slot, slide_schema, payload, result)

    def _check_divider_background(self, slide, slide_schema: SlideSchema,
                                  result: QAResult) -> None:
        """Verify section divider has the brand-blue background fill."""
        try:
            bg = slide.background
            fill_color = bg.fill.fore_color.rgb
            expected_hex = self.schema.design.divider_bg.lstrip("#")
            actual_hex = str(fill_color)
            if actual_hex.upper() != expected_hex.upper():
                result.issues.append(Issue(
                    severity="error",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name="",
                    category="divider_background",
                    message=(
                        f"Divider background color {actual_hex} != "
                        f"expected {expected_hex}"
                    ),
                ))
        except Exception:
            result.issues.append(Issue(
                severity="error",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name="",
                category="divider_background",
                message="Divider slide missing background fill",
            ))

    # ------------------------------------------------------------------
    # Per-slot checks
    # ------------------------------------------------------------------

    def _check_slot(self, slide, slot: DataSlot, slide_schema: SlideSchema,
                    payload: dict[str, Any], result: QAResult) -> None:
        """Dispatch validation for a single slot."""
        checkers = {
            SlotType.KPI_VALUE: self._check_kpi_slot,
            SlotType.TABLE: self._check_table_slot,
            SlotType.CHART: self._check_chart_slot,
            SlotType.TEXT: self._check_text_slot,
            SlotType.STATIC: self._check_text_slot,
            SlotType.SECTION_DIVIDER: self._check_text_slot,
        }
        checker = checkers.get(slot.slot_type)
        if checker:
            checker(slide, slot, slide_schema, payload, result)

    # -- KPI validation -------------------------------------------------

    def _check_kpi_slot(self, slide, slot: DataSlot,
                        slide_schema: SlideSchema,
                        payload: dict[str, Any],
                        result: QAResult) -> None:
        """Validate KPI slot: value present and formatted correctly."""
        value = payload.get(slot.data_key)
        all_text = _all_text_on_slide(slide)

        if _is_missing(value):
            # With missing data, N/A should appear
            if "N/A" not in all_text:
                result.issues.append(Issue(
                    severity="warning",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="kpi_missing_na",
                    message=(
                        f"KPI '{slot.data_key}' is missing but "
                        f"N/A not found on slide"
                    ),
                ))
            return

        # Check formatted value appears on slide
        if slot.format_rule:
            formatted = format_value(value, slot.format_rule.format_type)
            if formatted not in all_text:
                result.issues.append(Issue(
                    severity="error",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="kpi_value_missing",
                    message=(
                        f"Formatted KPI value '{formatted}' for "
                        f"'{slot.data_key}' not found on slide"
                    ),
                ))

        # Check label rendered
        if slot.label and slot.label not in all_text:
            result.issues.append(Issue(
                severity="warning",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name=slot.name,
                category="kpi_label_missing",
                message=f"KPI label '{slot.label}' not found on slide",
            ))

        # Check variance rendered with correct coloring
        if slot.variance_key:
            var_value = payload.get(slot.variance_key)
            if not _is_missing(var_value):
                self._check_variance_color(
                    slide, slot, slide_schema, var_value, result,
                )

    def _check_variance_color(self, slide, slot: DataSlot,
                              slide_schema: SlideSchema,
                              var_value: float,
                              result: QAResult) -> None:
        """Verify variance text uses correct positive/negative color."""
        expected_color = variance_color(var_value)
        expected_hex = expected_color.lstrip("#").upper()

        # Determine expected variance text
        if slot.format_rule and slot.format_rule.format_type == FormatType.POINTS_CHANGE:
            var_text = format_value(var_value, FormatType.POINTS_CHANGE)
        else:
            var_text = format_value(var_value, FormatType.VARIANCE_PERCENTAGE)

        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if var_text in run.text:
                        found = True
                        if run.font.color and run.font.color.rgb:
                            actual_hex = str(run.font.color.rgb).upper()
                            if actual_hex != expected_hex:
                                result.issues.append(Issue(
                                    severity="error",
                                    slide_index=slide_schema.index,
                                    slide_name=slide_schema.name,
                                    slot_name=slot.name,
                                    category="variance_color",
                                    message=(
                                        f"Variance '{var_text}' color "
                                        f"{actual_hex} != expected "
                                        f"{expected_hex}"
                                    ),
                                ))

        if not found:
            result.issues.append(Issue(
                severity="warning",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name=slot.name,
                category="variance_text_missing",
                message=(
                    f"Variance text '{var_text}' not found on slide"
                ),
            ))

    # -- Table validation -----------------------------------------------

    def _check_table_slot(self, slide, slot: DataSlot,
                          slide_schema: SlideSchema,
                          payload: dict[str, Any],
                          result: QAResult) -> None:
        """Validate table: row count, column headers, cell formatting."""
        rows_data = (
            payload.get(slot.row_data_key) if slot.row_data_key else None
        )

        tables = _table_shapes(slide)

        if not rows_data or not slot.columns:
            # No data — table should either not exist or be a placeholder
            return

        if not tables:
            result.issues.append(Issue(
                severity="error",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name=slot.name,
                category="table_missing",
                message="Table slot has data but no table shape on slide",
            ))
            return

        # Find the best matching table (by column count)
        table = None
        for ts in tables:
            t = ts.table
            if len(t.columns) == len(slot.columns):
                table = t
                break
        if table is None:
            table = tables[0].table

        # Check row count: header + data rows
        expected_rows = len(rows_data) + 1
        actual_rows = len(table.rows)
        if actual_rows != expected_rows:
            result.issues.append(Issue(
                severity="error",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name=slot.name,
                category="table_row_count",
                message=(
                    f"Table has {actual_rows} rows, expected "
                    f"{expected_rows} (1 header + {len(rows_data)} data)"
                ),
            ))

        # Check column count
        expected_cols = len(slot.columns)
        actual_cols = len(table.columns)
        if actual_cols != expected_cols:
            result.issues.append(Issue(
                severity="error",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name=slot.name,
                category="table_column_count",
                message=(
                    f"Table has {actual_cols} columns, "
                    f"expected {expected_cols}"
                ),
            ))
            return  # Skip header/cell checks if columns don't match

        # Check header text
        for col_idx, col_def in enumerate(slot.columns):
            header_text = table.cell(0, col_idx).text.strip()
            if header_text != col_def.header:
                result.issues.append(Issue(
                    severity="error",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="table_header",
                    message=(
                        f"Column {col_idx} header '{header_text}' != "
                        f"expected '{col_def.header}'"
                    ),
                ))

        # Check data cell formatting
        for row_idx, row_data in enumerate(rows_data):
            if row_idx + 1 >= len(table.rows):
                break
            for col_idx, col_def in enumerate(slot.columns):
                if col_idx >= len(table.columns):
                    break
                raw_val = row_data.get(col_def.data_key)
                cell_text = table.cell(row_idx + 1, col_idx).text.strip()

                if col_def.format_rule and not _is_missing(raw_val):
                    expected_text = format_value(
                        raw_val, col_def.format_rule.format_type,
                    )
                    if cell_text != expected_text:
                        result.issues.append(Issue(
                            severity="error",
                            slide_index=slide_schema.index,
                            slide_name=slide_schema.name,
                            slot_name=slot.name,
                            category="table_cell_format",
                            message=(
                                f"Cell [{row_idx+1},{col_idx}] "
                                f"'{cell_text}' != expected "
                                f"'{expected_text}' "
                                f"(format: {col_def.format_rule.format_type.value})"
                            ),
                        ))

        # Check variance coloring in table cells
        self._check_table_variance_colors(
            table, slot, slide_schema, rows_data, result,
        )

    def _check_table_variance_colors(self, table, slot: DataSlot,
                                     slide_schema: SlideSchema,
                                     rows_data: list[dict],
                                     result: QAResult) -> None:
        """Verify variance-colored cells in a table."""
        for col_idx, col_def in enumerate(slot.columns):
            if not col_def.format_rule:
                continue
            if col_def.format_rule.format_type not in (
                FormatType.VARIANCE_PERCENTAGE, FormatType.POINTS_CHANGE,
            ):
                continue

            for row_idx, row_data in enumerate(rows_data):
                if row_idx + 1 >= len(table.rows):
                    break
                if col_idx >= len(table.columns):
                    break
                raw_val = row_data.get(col_def.data_key)
                if _is_missing(raw_val):
                    continue

                expected_hex = variance_color(
                    raw_val,
                    col_def.format_rule.positive_color,
                    col_def.format_rule.negative_color,
                    col_def.format_rule.neutral_color,
                ).lstrip("#").upper()

                cell = table.cell(row_idx + 1, col_idx)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            actual_hex = str(run.font.color.rgb).upper()
                            if actual_hex != expected_hex:
                                result.issues.append(Issue(
                                    severity="error",
                                    slide_index=slide_schema.index,
                                    slide_name=slide_schema.name,
                                    slot_name=slot.name,
                                    category="table_variance_color",
                                    message=(
                                        f"Cell [{row_idx+1},{col_idx}] "
                                        f"variance color {actual_hex} != "
                                        f"expected {expected_hex} "
                                        f"(value={raw_val})"
                                    ),
                                ))

    # -- Chart validation -----------------------------------------------

    def _check_chart_slot(self, slide, slot: DataSlot,
                          slide_schema: SlideSchema,
                          payload: dict[str, Any],
                          result: QAResult) -> None:
        """Validate chart: type, series count, category count."""
        if not slot.chart_type or not slot.series:
            return

        charts = _chart_shapes(slide)
        if not charts:
            # Only an error if there was actual data to render
            has_data = any(
                payload.get(s.data_key) is not None for s in slot.series
            )
            if has_data:
                result.issues.append(Issue(
                    severity="error",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="chart_missing",
                    message="Chart slot has data but no chart shape on slide",
                ))
            return

        # Match chart shape to slot by chart type, then by position
        from pptx.enum.chart import XL_CHART_TYPE
        _type_map = {
            ChartType.COLUMN_CLUSTERED: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartType.LINE: XL_CHART_TYPE.LINE,
            ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
            ChartType.DOUGHNUT_EXPLODED: XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        }
        expected_type = _type_map.get(slot.chart_type)
        matched_shape = None
        if expected_type is not None:
            for cs in charts:
                if cs.chart.chart_type == expected_type:
                    matched_shape = cs
                    break
        if matched_shape is None:
            # Fallback: match by position proximity
            slot_left = Inches(slot.position.left)
            slot_top = Inches(slot.position.top)
            best = charts[0]
            best_dist = abs(best.left - slot_left) + abs(best.top - slot_top)
            for cs in charts[1:]:
                dist = abs(cs.left - slot_left) + abs(cs.top - slot_top)
                if dist < best_dist:
                    best = cs
                    best_dist = dist
            matched_shape = best
        chart = matched_shape.chart

        # Check chart type
        if expected_type and chart.chart_type != expected_type:
            result.issues.append(Issue(
                severity="error",
                slide_index=slide_schema.index,
                slide_name=slide_schema.name,
                slot_name=slot.name,
                category="chart_type",
                message=(
                    f"Chart type {chart.chart_type} != "
                    f"expected {expected_type}"
                ),
            ))

        # Check series count (non-doughnut)
        if slot.chart_type not in (
            ChartType.DOUGHNUT, ChartType.DOUGHNUT_EXPLODED,
        ):
            # Count series that have data in payload
            expected_series = sum(
                1 for s in slot.series
                if not _is_missing(payload.get(s.data_key))
                and payload.get(s.data_key)
            )
            if expected_series == 0:
                # If no data, builder may use zeros for all series
                expected_series = len(slot.series)
            actual_series = len(chart.series)
            if actual_series != expected_series:
                result.issues.append(Issue(
                    severity="warning",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="chart_series_count",
                    message=(
                        f"Chart has {actual_series} series, "
                        f"expected {expected_series}"
                    ),
                ))

        # Check categories length matches series data length
        if slot.categories_key:
            categories = payload.get(slot.categories_key)
            if categories and isinstance(categories, (list, tuple)):
                for s in slot.series:
                    series_data = payload.get(s.data_key)
                    if (
                        series_data
                        and isinstance(series_data, (list, tuple))
                        and len(series_data) != len(categories)
                    ):
                        result.issues.append(Issue(
                            severity="error",
                            slide_index=slide_schema.index,
                            slide_name=slide_schema.name,
                            slot_name=slot.name,
                            category="chart_data_length",
                            message=(
                                f"Series '{s.name}' has "
                                f"{len(series_data)} values but "
                                f"{len(categories)} categories"
                            ),
                        ))

    # -- Text validation ------------------------------------------------

    def _check_text_slot(self, slide, slot: DataSlot,
                         slide_schema: SlideSchema,
                         payload: dict[str, Any],
                         result: QAResult) -> None:
        """Validate text slot: content rendered on slide."""
        value = payload.get(slot.data_key)
        if _is_missing(value):
            return

        all_text = _all_text_on_slide(slide)

        if isinstance(value, list):
            for item in value:
                item_str = str(item)
                if item_str not in all_text:
                    result.issues.append(Issue(
                        severity="warning",
                        slide_index=slide_schema.index,
                        slide_name=slide_schema.name,
                        slot_name=slot.name,
                        category="text_content",
                        message=(
                            f"List item '{item_str}' not found on slide"
                        ),
                    ))
        elif isinstance(value, str) and value:
            if value not in all_text:
                result.issues.append(Issue(
                    severity="warning",
                    slide_index=slide_schema.index,
                    slide_name=slide_schema.name,
                    slot_name=slot.name,
                    category="text_content",
                    message=(
                        f"Text '{value[:60]}' not found on slide"
                    ),
                ))

    # ------------------------------------------------------------------
    # Utility
    # ------------------------------------------------------------------

    def _find_slide_for_key(self, data_key: str) -> str:
        """Find which slide a data_key belongs to."""
        for slide_schema in self.schema.slides:
            for slot in slide_schema.slots:
                keys = {slot.data_key}
                if slot.variance_key:
                    keys.add(slot.variance_key)
                if slot.row_data_key:
                    keys.add(slot.row_data_key)
                if slot.categories_key:
                    keys.add(slot.categories_key)
                for s in slot.series:
                    keys.add(s.data_key)
                if data_key in keys:
                    return slide_schema.name
        return ""


# ---------------------------------------------------------------------------
# Convenience function
# ---------------------------------------------------------------------------

def validate_presentation(schema: TemplateSchema, pptx_bytes: bytes,
                          payload: dict[str, Any]) -> QAResult:
    """One-shot convenience: validate a PPTX against its schema."""
    return QAValidator(schema).validate(pptx_bytes, payload)
