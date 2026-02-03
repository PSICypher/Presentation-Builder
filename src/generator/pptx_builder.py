"""PPTX builder engine — generates PowerPoint presentations from schema + data.

Consumes a TemplateSchema (slide definitions, data slots, design system) and a
data payload dict (keyed by data_key identifiers) to produce a fully rendered
.pptx file using python-pptx.

Usage::

    from src.generator.pptx_builder import PPTXBuilder
    from src.schema import build_monthly_report_schema

    schema = build_monthly_report_schema()
    builder = PPTXBuilder(schema)
    pptx_bytes = builder.build(payload)

    with open("report.pptx", "wb") as f:
        f.write(pptx_bytes)
"""

import io
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.schema.design_system import format_value, variance_color
from src.schema.models import (
    ChartType,
    DataSlot,
    DesignSystem,
    FontSpec,
    FormatRule,
    FormatType,
    Position,
    SlideSchema,
    SlideType,
    SlotType,
    TableColumn,
    TemplateSchema,
)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_CHART_TYPE_MAP = {
    ChartType.COLUMN_CLUSTERED: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE,
    ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
    ChartType.DOUGHNUT_EXPLODED: XL_CHART_TYPE.DOUGHNUT_EXPLODED,
}

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

_NA = "N/A"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' hex string to an RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(*bytes.fromhex(h))


def _is_missing(value: Any) -> bool:
    """Check if a value is None or NaN."""
    if value is None:
        return True
    if isinstance(value, float) and math.isnan(value):
        return True
    return False


def _format_slot_value(value: Any, format_rule: FormatRule | None) -> str:
    """Format a value using the slot's format rule, or fallback to str."""
    if _is_missing(value):
        return _NA
    if format_rule is None:
        return str(value) if value is not None else _NA
    return format_value(value, format_rule.format_type)


def _apply_font(run, font_spec: FontSpec | None, design: DesignSystem) -> None:
    """Apply a FontSpec to a python-pptx Run."""
    if font_spec is None:
        run.font.name = design.primary_font
        run.font.size = Pt(design.body_size_pt)
        return
    run.font.name = font_spec.name
    run.font.size = Pt(font_spec.size_pt)
    run.font.bold = font_spec.bold
    run.font.italic = font_spec.italic
    if font_spec.color:
        run.font.color.rgb = _hex_to_rgb(font_spec.color)


# ---------------------------------------------------------------------------
# PPTXBuilder
# ---------------------------------------------------------------------------

class PPTXBuilder:
    """Builds a PowerPoint presentation from a schema and data payload.

    Parameters
    ----------
    schema : TemplateSchema
        The full presentation schema including slide definitions, data slots,
        and design system.
    """

    def __init__(self, schema: TemplateSchema) -> None:
        self.schema = schema
        self.design = schema.design

    def build(self, payload: dict[str, Any]) -> bytes:
        """Build the PPTX and return it as bytes.

        Parameters
        ----------
        payload : dict[str, Any]
            Data payload keyed by data_key identifiers from the schema.
            Typically produced by ``DataMapper.map().payload``.

        Returns
        -------
        bytes
            The rendered .pptx file content.
        """
        prs = Presentation()
        prs.slide_width = Inches(self.schema.width_inches)
        prs.slide_height = Inches(self.schema.height_inches)

        for slide_schema in self.schema.slides:
            self._build_slide(prs, slide_schema, payload)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def build_to_file(self, payload: dict[str, Any], path: str | Path) -> None:
        """Build the PPTX and write it to a file path."""
        data = self.build(payload)
        Path(path).write_bytes(data)

    # ------------------------------------------------------------------
    # Slide builders
    # ------------------------------------------------------------------

    def _build_slide(self, prs: Presentation, slide_schema: SlideSchema,
                     payload: dict[str, Any]) -> None:
        """Create a slide and render all its slots."""
        # Use blank layout to avoid placeholder interference
        layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(layout)

        # Section divider gets a full-slide background fill
        if slide_schema.slide_type == SlideType.SECTION_DIVIDER:
            self._apply_divider_background(slide)

        for slot in slide_schema.slots:
            self._render_slot(slide, slot, payload)

    def _apply_divider_background(self, slide) -> None:
        """Apply brand-blue solid fill to the entire slide background."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(self.design.divider_bg)

    # ------------------------------------------------------------------
    # Slot renderers — dispatch by SlotType
    # ------------------------------------------------------------------

    def _render_slot(self, slide, slot: DataSlot, payload: dict[str, Any]) -> None:
        """Dispatch to the appropriate renderer based on slot type."""
        renderers = {
            SlotType.KPI_VALUE: self._render_kpi,
            SlotType.TABLE: self._render_table,
            SlotType.CHART: self._render_chart,
            SlotType.TEXT: self._render_text,
            SlotType.STATIC: self._render_text,
            SlotType.SECTION_DIVIDER: self._render_section_divider,
            SlotType.IMAGE: self._render_placeholder,
        }
        renderer = renderers.get(slot.slot_type, self._render_placeholder)
        renderer(slide, slot, payload)

    # ------------------------------------------------------------------
    # KPI rendering
    # ------------------------------------------------------------------

    def _render_kpi(self, slide, slot: DataSlot, payload: dict[str, Any]) -> None:
        """Render a KPI value with label and optional variance indicator."""
        pos = slot.position
        txbox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        # Label (above the number)
        if slot.label:
            p_label = tf.paragraphs[0]
            p_label.alignment = PP_ALIGN.CENTER
            run_label = p_label.add_run()
            run_label.text = slot.label
            label_font = FontSpec(
                name=self.design.primary_font,
                size_pt=self.design.kpi_label_size_pt,
                color=self.design.dark_grey,
            )
            _apply_font(run_label, label_font, self.design)

        # Main value
        value = payload.get(slot.data_key)
        formatted = _format_slot_value(value, slot.format_rule)

        p_value = tf.add_paragraph()
        p_value.alignment = PP_ALIGN.CENTER
        run_value = p_value.add_run()
        run_value.text = formatted
        _apply_font(run_value, slot.font, self.design)

        # Variance indicator
        if slot.variance_key:
            var_value = payload.get(slot.variance_key)
            if not _is_missing(var_value):
                var_text = format_value(var_value, FormatType.VARIANCE_PERCENTAGE)
                var_color = variance_color(var_value)
                if slot.format_rule and slot.format_rule.format_type == FormatType.POINTS_CHANGE:
                    var_text = format_value(var_value, FormatType.POINTS_CHANGE)

                p_var = tf.add_paragraph()
                p_var.alignment = PP_ALIGN.CENTER
                run_var = p_var.add_run()
                run_var.text = var_text
                run_var.font.name = self.design.primary_font
                run_var.font.size = Pt(self.design.caption_size_pt)
                run_var.font.color.rgb = _hex_to_rgb(var_color)

    # ------------------------------------------------------------------
    # Table rendering
    # ------------------------------------------------------------------

    def _render_table(self, slide, slot: DataSlot, payload: dict[str, Any]) -> None:
        """Render a data table from row_data_key and column definitions."""
        rows_data = payload.get(slot.row_data_key) if slot.row_data_key else None
        if not rows_data or not slot.columns:
            # Render empty table placeholder
            self._render_text(slide, slot, payload)
            return

        num_rows = len(rows_data) + 1  # +1 for header
        num_cols = len(slot.columns)
        pos = slot.position

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        table = table_shape.table

        # Set column widths
        for col_idx, col_def in enumerate(slot.columns):
            if col_def.width_inches is not None:
                table.columns[col_idx].width = Inches(col_def.width_inches)

        # Header row
        for col_idx, col_def in enumerate(slot.columns):
            cell = table.cell(0, col_idx)
            cell.text = col_def.header
            self._style_table_cell(
                cell, is_header=True, alignment=col_def.alignment,
                font_spec=col_def.font,
            )

        # Data rows
        for row_idx, row_data in enumerate(rows_data):
            for col_idx, col_def in enumerate(slot.columns):
                cell = table.cell(row_idx + 1, col_idx)
                raw_val = row_data.get(col_def.data_key)
                formatted = _format_slot_value(raw_val, col_def.format_rule)
                cell.text = formatted

                # Apply variance coloring for variance columns
                color_override = None
                if col_def.format_rule and col_def.format_rule.format_type in (
                    FormatType.VARIANCE_PERCENTAGE, FormatType.POINTS_CHANGE,
                ):
                    if not _is_missing(raw_val):
                        color_override = variance_color(
                            raw_val,
                            col_def.format_rule.positive_color,
                            col_def.format_rule.negative_color,
                            col_def.format_rule.neutral_color,
                        )

                self._style_table_cell(
                    cell, is_header=False, alignment=col_def.alignment,
                    color_override=color_override,
                )

    def _style_table_cell(self, cell, is_header: bool = False,
                          alignment: str = "left",
                          font_spec: FontSpec | None = None,
                          color_override: str | None = None) -> None:
        """Apply styling to a table cell."""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = _ALIGN_MAP.get(alignment, PP_ALIGN.LEFT)
            for run in paragraph.runs:
                run.font.name = self.design.primary_font
                if is_header:
                    run.font.size = Pt(11.0)
                    run.font.bold = True
                    run.font.color.rgb = _hex_to_rgb(self.design.white)
                else:
                    run.font.size = Pt(11.0)
                    run.font.bold = False
                    if color_override:
                        run.font.color.rgb = _hex_to_rgb(color_override)
                    else:
                        run.font.color.rgb = _hex_to_rgb(self.design.dark_text)

        # Header row background
        if is_header:
            from pptx.oxml.ns import qn
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn("a:solidFill"), {})
            srgbClr = solidFill.makeelement(
                qn("a:srgbClr"),
                {"val": self.design.dark_blue.lstrip("#")},
            )
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

    # ------------------------------------------------------------------
    # Chart rendering
    # ------------------------------------------------------------------

    def _render_chart(self, slide, slot: DataSlot, payload: dict[str, Any]) -> None:
        """Render a chart from series data and categories."""
        if not slot.chart_type or not slot.series:
            self._render_placeholder(slide, slot, payload)
            return

        categories = payload.get(slot.categories_key) if slot.categories_key else None
        xl_chart_type = _CHART_TYPE_MAP.get(slot.chart_type)
        if xl_chart_type is None:
            self._render_placeholder(slide, slot, payload)
            return

        chart_data = CategoryChartData()

        # Set categories
        if categories and not _is_missing(categories):
            chart_data.categories = [str(c) for c in categories]
        else:
            # For doughnut charts without explicit categories, use series names
            if slot.chart_type in (ChartType.DOUGHNUT, ChartType.DOUGHNUT_EXPLODED):
                chart_data.categories = [s.name for s in slot.series]
                # Doughnut: single series with values from each named segment
                values = []
                for s in slot.series:
                    val = payload.get(s.data_key)
                    values.append(val if not _is_missing(val) else 0)
                chart_data.add_series("Data", tuple(values))

                pos = slot.position
                chart_frame = slide.shapes.add_chart(
                    xl_chart_type,
                    Inches(pos.left), Inches(pos.top),
                    Inches(pos.width), Inches(pos.height),
                    chart_data,
                )
                chart = chart_frame.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

                # Apply series colors
                plot = chart.plots[0]
                for idx, s in enumerate(slot.series):
                    if s.color:
                        point = plot.series[0].points[idx]
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = _hex_to_rgb(s.color)
                return

        # Standard category charts (column, line)
        for s in slot.series:
            series_values = payload.get(s.data_key)
            if _is_missing(series_values) or not series_values:
                # Use zeros if no data
                if categories:
                    series_values = [0] * len(categories)
                else:
                    continue
            chart_data.add_series(s.name, tuple(series_values))

        pos = slot.position
        chart_frame = slide.shapes.add_chart(
            xl_chart_type,
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
            chart_data,
        )
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        # Apply series colors
        plot = chart.plots[0]
        for idx, s in enumerate(slot.series):
            if s.color and idx < len(plot.series):
                series_obj = plot.series[idx]
                series_obj.format.fill.solid()
                series_obj.format.fill.fore_color.rgb = _hex_to_rgb(s.color)

    # ------------------------------------------------------------------
    # Text rendering
    # ------------------------------------------------------------------

    def _render_text(self, slide, slot: DataSlot, payload: dict[str, Any]) -> None:
        """Render a text box with content from the payload."""
        pos = slot.position
        txbox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        value = payload.get(slot.data_key)
        if _is_missing(value):
            text = ""
        elif isinstance(value, list):
            # Handle list of items (e.g., TOC items, bullet points)
            text = "\n".join(str(item) for item in value)
        else:
            text = str(value)

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _apply_font(run, slot.font, self.design)

    # ------------------------------------------------------------------
    # Section divider rendering
    # ------------------------------------------------------------------

    def _render_section_divider(self, slide, slot: DataSlot,
                                payload: dict[str, Any]) -> None:
        """Render section divider title text (white on brand-blue background)."""
        pos = slot.position
        txbox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Use the slide title from the schema as divider text
        value = payload.get(slot.data_key)
        text = str(value) if not _is_missing(value) else slot.name

        run = tf.paragraphs[0].add_run()
        run.text = text
        _apply_font(run, slot.font, self.design)

        # Vertically center the text
        txbox.text_frame.paragraphs[0].space_before = Pt(0)
        txbox.text_frame.paragraphs[0].space_after = Pt(0)

    # ------------------------------------------------------------------
    # Placeholder for unsupported types
    # ------------------------------------------------------------------

    def _render_placeholder(self, slide, slot: DataSlot,
                            payload: dict[str, Any]) -> None:
        """Render a placeholder shape for unsupported slot types."""
        pos = slot.position
        txbox = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top),
            Inches(pos.width), Inches(pos.height),
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"[{slot.slot_type.value}: {slot.name}]"
        run.font.name = self.design.primary_font
        run.font.size = Pt(self.design.caption_size_pt)
        run.font.color.rgb = _hex_to_rgb(self.design.light_gray)


# ---------------------------------------------------------------------------
# Convenience function
# ---------------------------------------------------------------------------

def build_presentation(schema: TemplateSchema, payload: dict[str, Any]) -> bytes:
    """One-shot convenience: build a PPTX from schema + payload."""
    return PPTXBuilder(schema).build(payload)
