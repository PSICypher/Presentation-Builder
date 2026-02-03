"""Chart generation module — renders chart shapes on PowerPoint slides.

Converts DataSlot chart specifications and payload data into python-pptx
chart shapes with proper styling, colors, and positioning.

Supported chart types:
    COLUMN_CLUSTERED — Grouped column chart (multi-series side-by-side)
    LINE             — Line chart (multi-series)
    DOUGHNUT         — Donut chart (single series, multiple slices)
    DOUGHNUT_EXPLODED — Exploded donut chart (separated slices)

Usage:
    from src.generator.charts import add_chart

    added = add_chart(slide, slot, payload, design)
"""

from __future__ import annotations

import math
from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from src.schema.models import (
    ChartType,
    DataSlot,
    DesignSystem,
    SlideSchema,
    SlotType,
)


# ---------------------------------------------------------------------------
# Chart type mapping
# ---------------------------------------------------------------------------

_CHART_TYPE_MAP: dict[ChartType, int] = {
    ChartType.COLUMN_CLUSTERED: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE,
    ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
    ChartType.DOUGHNUT_EXPLODED: XL_CHART_TYPE.DOUGHNUT_EXPLODED,
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string (#RRGGBB) to an RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _safe_value(value: Any) -> float:
    """Coerce a value to a safe float for chart data.  None/NaN/inf → 0."""
    if value is None:
        return 0.0
    if isinstance(value, (int, float)):
        if math.isnan(value) or math.isinf(value):
            return 0.0
        return float(value)
    return 0.0


def _is_doughnut(chart_type: ChartType) -> bool:
    """Check if the chart type is a doughnut variant."""
    return chart_type in (ChartType.DOUGHNUT, ChartType.DOUGHNUT_EXPLODED)


# ---------------------------------------------------------------------------
# Chart data builders
# ---------------------------------------------------------------------------

def _build_category_chart_data(
    slot: DataSlot,
    payload: dict[str, Any],
) -> CategoryChartData | None:
    """Build chart data for category-based charts (column, line).

    Returns None if categories are missing or no series definitions exist.
    """
    if not slot.categories_key:
        return None

    categories = payload.get(slot.categories_key)
    if not categories:
        return None

    if not slot.series:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series_def in slot.series:
        values = payload.get(series_def.data_key)
        if values is None:
            values = [0.0] * len(categories)
        else:
            values = list(values)
            if len(values) < len(categories):
                values += [0.0] * (len(categories) - len(values))
            elif len(values) > len(categories):
                values = values[: len(categories)]

        safe_values = tuple(_safe_value(v) for v in values)
        chart_data.add_series(series_def.name, safe_values)

    return chart_data


def _build_doughnut_chart_data(
    slot: DataSlot,
    payload: dict[str, Any],
) -> CategoryChartData | None:
    """Build chart data for doughnut charts.

    Each series definition becomes a single slice in the doughnut.
    Returns None if no series exist or all values resolve to zero.
    """
    if not slot.series:
        return None

    categories = [s.name for s in slot.series]
    values = [_safe_value(payload.get(s.data_key)) for s in slot.series]

    if all(v == 0.0 for v in values):
        return None

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Data", tuple(values))

    return chart_data


# ---------------------------------------------------------------------------
# Chart styling
# ---------------------------------------------------------------------------

def _apply_series_colors(chart, slot: DataSlot) -> None:
    """Apply per-series colors from the slot definition."""
    plot = chart.plots[0]

    if _is_doughnut(slot.chart_type):
        # Doughnut: colors go on individual points of the single series
        if plot.series and slot.series:
            series_obj = plot.series[0]
            for idx, series_def in enumerate(slot.series):
                if series_def.color and idx < len(list(series_obj.points)):
                    point = series_obj.points[idx]
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = _hex_to_rgb(series_def.color)
    else:
        # Column/Line: colors go on each series object
        for idx, series_def in enumerate(slot.series):
            if idx < len(plot.series) and series_def.color:
                plot_series = plot.series[idx]
                rgb = _hex_to_rgb(series_def.color)
                if slot.chart_type == ChartType.LINE:
                    plot_series.format.line.color.rgb = rgb
                else:
                    plot_series.format.fill.solid()
                    plot_series.format.fill.fore_color.rgb = rgb


def _apply_chart_style(chart, slot: DataSlot, design: DesignSystem) -> None:
    """Apply general styling: font, legend visibility and position."""
    chart.font.name = design.primary_font
    chart.font.size = Pt(design.caption_size_pt)

    if _is_doughnut(slot.chart_type):
        chart.has_legend = False
    elif len(slot.series) > 1:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.name = design.primary_font
        chart.legend.font.size = Pt(design.caption_size_pt)
    else:
        chart.has_legend = False


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def add_chart(
    slide,
    slot: DataSlot,
    payload: dict[str, Any],
    design: DesignSystem,
) -> bool:
    """Add a chart shape to a slide based on slot configuration.

    Args:
        slide: python-pptx Slide object.
        slot: DataSlot with chart_type, series, categories_key, and position.
        payload: Data payload dict with values keyed by data_key.
        design: DesignSystem for styling.

    Returns:
        True if the chart was added, False if skipped due to missing data.

    Raises:
        ValueError: If the slot is not a CHART type or has no chart_type.
    """
    if slot.slot_type != SlotType.CHART:
        raise ValueError(f"Slot '{slot.name}' is not a CHART type")
    if slot.chart_type is None:
        raise ValueError(f"Slot '{slot.name}' has no chart_type")

    if _is_doughnut(slot.chart_type):
        chart_data = _build_doughnut_chart_data(slot, payload)
    else:
        chart_data = _build_category_chart_data(slot, payload)

    if chart_data is None:
        return False

    xl_chart_type = _CHART_TYPE_MAP[slot.chart_type]

    pos = slot.position
    graphic_frame = slide.shapes.add_chart(
        xl_chart_type,
        Inches(pos.left),
        Inches(pos.top),
        Inches(pos.width),
        Inches(pos.height),
        chart_data,
    )
    chart = graphic_frame.chart

    _apply_series_colors(chart, slot)
    _apply_chart_style(chart, slot, design)

    return True


def add_slide_charts(
    slide,
    slide_schema: SlideSchema,
    payload: dict[str, Any],
    design: DesignSystem,
) -> list[str]:
    """Add all chart-type slots on a slide.

    Returns the list of slot names that were successfully rendered.
    """
    added: list[str] = []
    for slot in slide_schema.slots:
        if slot.slot_type == SlotType.CHART:
            if add_chart(slide, slot, payload, design):
                added.append(slot.name)
    return added
