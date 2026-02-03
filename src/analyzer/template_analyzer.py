"""Template Analyzer - Extracts structure, styling, and patterns from PPTX files.

Uses python-pptx to deeply inspect PowerPoint files and produce a machine-readable
analysis (JSON/YAML) that informs the template schema model and extraction engine.
"""

from collections import Counter
from pathlib import Path
from typing import Any

import yaml
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def _emu_to_inches(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / EMU_PER_INCH, 3)


def _emu_to_pt(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / EMU_PER_PT, 1)


def _safe_get(fn):
    """Call fn, return None on error."""
    try:
        return fn()
    except Exception:
        return None


class TemplateAnalyzer:
    """Analyses a PPTX file and extracts its complete structure."""

    def __init__(self, path: str | Path):
        self.path = Path(path)
        self.prs = Presentation(str(self.path))
        self._fonts: Counter = Counter()
        self._font_sizes: Counter = Counter()
        self._colors: Counter = Counter()
        self._shape_types: Counter = Counter()
        self._layout_usage: Counter = Counter()

    def analyze(self) -> dict[str, Any]:
        """Run full analysis and return structured result."""
        result = {
            "source_file": self.path.name,
            "dimensions": self._extract_dimensions(),
            "theme": self._extract_theme(),
            "slide_masters": self._extract_masters(),
            "slides": self._extract_slides(),
            "summary": {},
        }

        result["summary"] = {
            "slide_count": len(result["slides"]),
            "fonts": dict(self._fonts.most_common()),
            "font_sizes_pt": dict(self._font_sizes.most_common()),
            "colors_hex": dict(self._colors.most_common(20)),
            "shape_types": dict(self._shape_types.most_common()),
            "layout_usage": dict(self._layout_usage.most_common()),
            "total_tables": sum(1 for s in result["slides"] for sh in s["shapes"] if sh.get("table")),
            "total_charts": sum(1 for s in result["slides"] for sh in s["shapes"] if sh.get("chart")),
            "total_images": sum(1 for s in result["slides"] for sh in s["shapes"] if sh.get("is_picture")),
            "total_groups": sum(1 for s in result["slides"] for sh in s["shapes"] if sh.get("is_group")),
        }

        return result

    def _extract_dimensions(self) -> dict:
        return {
            "width_inches": _emu_to_inches(self.prs.slide_width),
            "height_inches": _emu_to_inches(self.prs.slide_height),
            "width_emu": self.prs.slide_width,
            "height_emu": self.prs.slide_height,
        }

    def _extract_theme(self) -> dict:
        """Extract theme colors and fonts from the slide master's theme part."""
        themes = []
        for master in self.prs.slide_masters:
            for rel in master.part.rels.values():
                if "theme" not in rel.reltype:
                    continue
                theme_xml = rel.target_part.blob
                root = etree.fromstring(theme_xml)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

                theme = {}

                # Color scheme
                clr_scheme = root.find(".//a:clrScheme", ns)
                if clr_scheme is not None:
                    theme["color_scheme_name"] = clr_scheme.get("name", "")
                    colors = {}
                    for child in clr_scheme:
                        tag = child.tag.split("}")[-1]
                        for sub in child:
                            val = sub.get("val") or sub.get("lastClr")
                            if val:
                                colors[tag] = val
                    theme["theme_colors"] = colors

                # Font scheme
                font_scheme = root.find(".//a:fontScheme", ns)
                if font_scheme is not None:
                    theme["font_scheme_name"] = font_scheme.get("name", "")
                    major = font_scheme.find(".//a:majorFont/a:latin", ns)
                    minor = font_scheme.find(".//a:minorFont/a:latin", ns)
                    theme["major_font"] = major.get("typeface", "") if major is not None else ""
                    theme["minor_font"] = minor.get("typeface", "") if minor is not None else ""

                themes.append(theme)
                break

        return themes[0] if len(themes) == 1 else {"masters": themes}

    def _extract_masters(self) -> list[dict]:
        """Extract slide master and layout information."""
        masters = []
        for master in self.prs.slide_masters:
            m = {"layouts": []}
            for layout in master.slide_layouts:
                lay = {"name": layout.name, "placeholders": []}
                try:
                    for ph in layout.placeholders:
                        try:
                            lay["placeholders"].append({
                                "idx": ph.placeholder_format.idx,
                                "type": str(ph.placeholder_format.type),
                                "name": ph.name,
                                "position": {
                                    "left": _emu_to_inches(ph.left),
                                    "top": _emu_to_inches(ph.top),
                                    "width": _emu_to_inches(ph.width),
                                    "height": _emu_to_inches(ph.height),
                                },
                            })
                        except Exception:
                            pass
                except Exception:
                    pass
                m["layouts"].append(lay)
            masters.append(m)
        return masters

    def _extract_run_font(self, run) -> dict:
        """Extract font properties from a text run."""
        info = {}
        name = _safe_get(lambda: run.font.name)
        if name:
            info["name"] = name
            self._fonts[name] += 1

        size = _safe_get(lambda: run.font.size)
        if size:
            pt = _emu_to_pt(size)
            info["size_pt"] = pt
            self._font_sizes[str(pt)] += 1

        bold = _safe_get(lambda: run.font.bold)
        if bold:
            info["bold"] = True

        italic = _safe_get(lambda: run.font.italic)
        if italic:
            info["italic"] = True

        rgb = _safe_get(lambda: str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None)
        if rgb:
            info["color"] = rgb
            self._colors[rgb] += 1

        theme_clr = _safe_get(lambda: str(run.font.color.theme_color) if run.font.color and run.font.color.theme_color else None)
        if theme_clr:
            info["theme_color"] = theme_clr

        return info

    def _extract_shape_colors(self, shape) -> list[str]:
        """Extract fill/stroke colors from shape XML."""
        colors = []
        try:
            for elem in shape._element.iter():
                if elem.tag.endswith("}srgbClr"):
                    val = elem.get("val")
                    if val:
                        colors.append(val)
                        self._colors[val] += 1
        except Exception:
            pass
        return colors

    def _extract_text_content(self, shape) -> list[dict]:
        """Extract paragraphs and runs from a text frame."""
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            p = {"text": para.text, "runs": []}
            for run in para.runs:
                r = {"text": run.text}
                font = self._extract_run_font(run)
                if font:
                    r["font"] = font
                p["runs"].append(r)
            paragraphs.append(p)
        return paragraphs

    def _extract_table(self, shape) -> dict:
        """Extract table structure."""
        table = shape.table
        rows = list(table.rows)
        cols = list(table.columns)

        col_widths = [_emu_to_inches(c.width) for c in cols]

        # Extract header row content and formatting
        headers = []
        if rows:
            for cell in rows[0].cells:
                cell_info = {"text": cell.text.replace("\x0b", " ")}
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        font = self._extract_run_font(run)
                        if font:
                            cell_info["font"] = font
                            break
                    break
                headers.append(cell_info)

        return {
            "rows": len(rows),
            "cols": len(cols),
            "col_widths_inches": col_widths,
            "headers": headers,
        }

    def _extract_chart(self, shape) -> dict:
        """Extract chart metadata."""
        chart = shape.chart
        info = {"chart_type": str(chart.chart_type)}
        try:
            info["series_count"] = len(list(chart.series))
        except Exception:
            pass
        return info

    def _extract_shapes(self, slide) -> list[dict]:
        """Extract all shapes from a slide."""
        shapes = []
        for shape in slide.shapes:
            stype = str(shape.shape_type)
            self._shape_types[stype] += 1

            sh: dict[str, Any] = {
                "name": shape.name,
                "shape_type": stype,
                "position": {
                    "left": _emu_to_inches(shape.left),
                    "top": _emu_to_inches(shape.top),
                    "width": _emu_to_inches(shape.width),
                    "height": _emu_to_inches(shape.height),
                },
            }

            # Fill/stroke colors
            fill_colors = self._extract_shape_colors(shape)
            if fill_colors:
                sh["fill_colors"] = list(set(fill_colors))

            # Placeholder info
            try:
                if shape.placeholder_format is not None:
                    sh["placeholder_idx"] = shape.placeholder_format.idx
                    sh["placeholder_type"] = str(shape.placeholder_format.type)
            except Exception:
                pass

            # Text content
            if shape.has_text_frame:
                sh["text"] = self._extract_text_content(shape)

            # Table
            if shape.has_table:
                sh["table"] = self._extract_table(shape)

            # Chart
            if shape.has_chart:
                try:
                    sh["chart"] = self._extract_chart(shape)
                except Exception:
                    sh["chart"] = {"error": "could not parse"}

            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sh["is_picture"] = True

            # Group
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                sh["is_group"] = True
                sh["group_child_count"] = len(shape.shapes)

            shapes.append(sh)
        return shapes

    def _extract_slides(self) -> list[dict]:
        """Extract all slides with their shapes."""
        slides = []
        for i, slide in enumerate(self.prs.slides):
            layout_name = slide.slide_layout.name
            self._layout_usage[layout_name] += 1

            s = {
                "index": i,
                "layout": layout_name,
                "shapes": self._extract_shapes(slide),
            }
            slides.append(s)
        return slides


def analyze_template(path: str | Path) -> dict[str, Any]:
    """Convenience function: analyze a single PPTX file."""
    analyzer = TemplateAnalyzer(path)
    return analyzer.analyze()


def analyze_all_templates(template_dir: str | Path) -> list[dict[str, Any]]:
    """Analyze all .pptx files in a directory."""
    template_dir = Path(template_dir)
    results = []
    for pptx_file in sorted(template_dir.glob("*.pptx")):
        results.append(analyze_template(pptx_file))
    return results


def save_analysis(analysis: dict | list, output_path: str | Path, fmt: str = "json"):
    """Save analysis results to JSON or YAML."""
    import json

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if fmt == "yaml":
        with open(output_path, "w") as f:
            yaml.dump(analysis, f, default_flow_style=False, sort_keys=False, allow_unicode=True)
    else:
        with open(output_path, "w") as f:
            json.dump(analysis, f, indent=2, default=str)


if __name__ == "__main__":
    import argparse
    import json

    parser = argparse.ArgumentParser(description="Analyze PPTX template structure")
    parser.add_argument("path", help="Path to PPTX file or directory of PPTX files")
    parser.add_argument("-o", "--output", help="Output file path (JSON or YAML)")
    parser.add_argument("-f", "--format", choices=["json", "yaml"], default="json")
    parser.add_argument("--summary-only", action="store_true", help="Print summary only")
    args = parser.parse_args()

    target = Path(args.path)
    if target.is_dir():
        results = analyze_all_templates(target)
    else:
        results = [analyze_template(target)]

    for r in results:
        print(f"\n{'='*60}")
        print(f"  {r['source_file']}")
        print(f"{'='*60}")
        s = r["summary"]
        print(f"  Slides: {s['slide_count']}")
        print(f"  Dimensions: {r['dimensions']['width_inches']}\" x {r['dimensions']['height_inches']}\"")
        print(f"  Fonts: {s['fonts']}")
        print(f"  Top colors: {dict(list(s['colors_hex'].items())[:5])}")
        print(f"  Layouts: {s['layout_usage']}")
        print(f"  Tables: {s['total_tables']}, Charts: {s['total_charts']}, Images: {s['total_images']}")

    if args.output:
        data = results[0] if len(results) == 1 else results
        save_analysis(data, args.output, args.format)
        print(f"\n✓ Analysis saved to {args.output}")
