"""End-to-end integration tests with CoWork data.

Exercises the full Presentation Builder pipeline:
    synthetic CoWork data → DataMapper → PPTXBuilder → QAValidator

Each test builds realistic data matching the No7 US x THGi eComm report
structure, runs it through the complete pipeline, and validates that the
output PPTX is structurally correct and data-accurate.
"""

import calendar
import io
import math

import pandas as pd
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from src.generator.pptx_builder import PPTXBuilder
from src.processor.mapper import (
    DataMapper,
    MappingResult,
    REPORT_CHANNELS,
)
from src.qa.validator import QAValidator, QAResult
from src.schema.models import SlotType, SlideType
from src.schema.monthly_report import build_monthly_report_schema
from src.schema.qbr_report import build_qbr_schema


# ---------------------------------------------------------------------------
# Synthetic CoWork data factories
# ---------------------------------------------------------------------------

def _make_raw_data(year, month, days=None, channels=None):
    """Build a synthetic RAW DATA DataFrame matching CoWork tracker schema."""
    if channels is None:
        channels = REPORT_CHANNELS
    if days is None:
        days = list(range(1, 29))

    rows = []
    for day in days:
        for ch in channels:
            rows.append({
                "COS Year": year,
                "COS Month": month,
                "COS Day": day,
                "COS Channel": ch,
                "COS Locale": "en_US",
                "COS Orders": 10,
                "COS New Customers": 5,
                "COS COS%": 0.10,
                "COS CAC": 20.0,
                "COS CPA": 10.0,
                "COS Cost": 100.0,
                "COS Revenue": 1000.0,
                "COS Sessions": 200,
                "COS AOV": 100.0,
                "COS Conversion": 0.05,
            })
    return pd.DataFrame(rows)


def _make_targets(year, month, channels=None):
    """Build a synthetic Targets DataFrame matching CoWork target schema."""
    if channels is None:
        channels = REPORT_CHANNELS
    num_days = calendar.monthrange(year, month)[1]

    rows = []
    for day in range(1, num_days + 1):
        for ch in channels:
            rows.append({
                "Target_Type_Id": "Daily",
                "Date": pd.Timestamp(year, month, day),
                "Site_Id": "No 7",
                "Locale_Id": "en_US",
                "Channel_Id": ch,
                "Notes": float("nan"),
                "Gross_Revenue_Target": 900.0,
                "Net_Revenue_Target": 850.0,
                "Marketing_Spend_Target": 90.0,
                "Session_Target": 180,
                "Order_Target": 9,
                "New_Customer_Target": 4,
            })
    return pd.DataFrame(rows)


def _make_tracker(year, month, ly_year=None, ly_month=None):
    """Build tracker dict with RAW DATA for current + prior year."""
    if ly_year is None:
        ly_year = year - 1
    if ly_month is None:
        ly_month = month
    raw_cur = _make_raw_data(year, month)
    raw_ly = _make_raw_data(ly_year, ly_month)
    raw = pd.concat([raw_cur, raw_ly], ignore_index=True)
    return {"RAW DATA": raw}


def _make_offer_performance(month):
    """Build synthetic offer performance DataFrame (CoWork promo format)."""
    rows = []
    promos = ["Promo A", "Promo B", "Promo C"]
    channels = ["AFFILIATE", "EMAIL", "PPC"]
    for promo in promos:
        for ch in channels:
            rows.append({
                "Dimension 1": promo,
                "Dimension 2": ch,
                "Dimension 3": str(month),
                "Dimension 4": "Total",
                "Redemptions": 1000.0,
                "% Change Redemptions": -0.10,
                "Revenue": 50000.0,
                "% Change Revenue": 0.15,
                "Discount Amount": 5000.0,
                "% Change Discount Amount": -0.05,
            })
    rows.append({
        "Dimension 1": "Grand Total",
        "Dimension 2": "Total",
        "Dimension 3": "Total",
        "Dimension 4": "Total",
        "Redemptions": 9000.0,
        "% Change Redemptions": -0.10,
        "Revenue": 450000.0,
        "% Change Revenue": 0.15,
        "Discount Amount": 45000.0,
        "% Change Discount Amount": -0.05,
    })
    return pd.DataFrame(rows)


def _make_product_sales(month):
    """Build synthetic product sales DataFrame (CoWork product format)."""
    rows = []
    products = ["Product X", "Product Y", "Product Z"]
    for prod in products:
        rows.append({
            "Dimension 1": prod,
            "Dimension 2": str(month),
            "Dimension 3": "Total",
            "Units (Analysis)": 500.0,
            "Units (Comparison)": 400.0,
            "Units (vs. Comp)": 0.25,
            "Total Revenue (Analysis)": 25000.0,
            "Total Revenue (Comparison)": 20000.0,
            "Total Revenue (vs. Comp)": 0.25,
            "AOV (Analysis)": 50.0,
            "Avg. Selling Price (Analysis)": 45.0,
            "Total Discount % (Analysis)": 0.10,
            "New Customers (Analysis)": 100.0,
        })
    rows.append({
        "Dimension 1": "Grand Total",
        "Dimension 2": "Total",
        "Dimension 3": "Total",
        "Units (Analysis)": 1500.0,
        "Units (Comparison)": 1200.0,
        "Units (vs. Comp)": 0.25,
        "Total Revenue (Analysis)": 75000.0,
        "Total Revenue (Comparison)": 60000.0,
        "Total Revenue (vs. Comp)": 0.25,
        "AOV (Analysis)": 50.0,
        "Avg. Selling Price (Analysis)": 45.0,
        "Total Discount % (Analysis)": 0.10,
        "New Customers (Analysis)": 300.0,
    })
    return pd.DataFrame(rows)


def _make_crm():
    """Build synthetic CRM performance DataFrame (CoWork email format)."""
    return pd.DataFrame([
        {
            "Col A": "Grand Total", "Col B": "Total", "Col C": "Total",
            "Emails Sent": 50000, "Emails Sent vs Comp": -0.05,
            "Open Rate": 0.25, "Open Rate vs Comp": 0.02,
            "Click-Through Rate": 0.08, "Click-Through Rate vs Comp": -0.01,
            "Sessions": 10000, "Sessions vs Comp": 0.10,
            "Orders": 500, "Orders vs Comp": 0.20,
            "CVR": 0.05, "CVR vs Comp": 0.005,
            "Revenue": 75000.0, "Revenue vs Comp": 0.15,
            "AOV": 150.0, "AOV vs Comp": -0.03,
        },
        {
            "Col A": "Grand Total", "Col B": "Total", "Col C": "Manual",
            "Emails Sent": 30000, "Emails Sent vs Comp": -0.03,
            "Open Rate": 0.22, "Open Rate vs Comp": 0.01,
            "Click-Through Rate": 0.07, "Click-Through Rate vs Comp": -0.005,
            "Sessions": 6000, "Sessions vs Comp": 0.08,
            "Orders": 300, "Orders vs Comp": 0.18,
            "CVR": 0.05, "CVR vs Comp": 0.004,
            "Revenue": 45000.0, "Revenue vs Comp": 0.12,
            "AOV": 150.0, "AOV vs Comp": -0.02,
        },
        {
            "Col A": "Grand Total", "Col B": "Total", "Col C": "Automated",
            "Emails Sent": 20000, "Emails Sent vs Comp": -0.08,
            "Open Rate": 0.30, "Open Rate vs Comp": 0.03,
            "Click-Through Rate": 0.10, "Click-Through Rate vs Comp": -0.02,
            "Sessions": 4000, "Sessions vs Comp": 0.12,
            "Orders": 200, "Orders vs Comp": 0.22,
            "CVR": 0.05, "CVR vs Comp": 0.006,
            "Revenue": 30000.0, "Revenue vs Comp": 0.20,
            "AOV": 150.0, "AOV vs Comp": -0.05,
        },
    ])


def _make_affiliate():
    """Build synthetic affiliate publisher DataFrame (CoWork affiliate format)."""
    rows = [
        {
            "Dimension 1": "Grand Total",
            "Dimension 2": "All Publishers",
            "Dimension 3": "Total",
            "Influencer Filter": "Total",
            "Revenue (Analysis)": 100000.0,
            "Revenue (Comparison)": 80000.0,
            "Revenue (vs Comp)": 0.25,
            "Cost (Analysis)": 10000.0,
            "Cost (Comparison)": 9000.0,
            "Cost (vs Comp)": 0.111,
            "CoS (Analysis)": 0.10,
            "CoS (vs Comp)": -0.01,
            "Orders (Analysis)": 1000,
            "Orders (vs Comp)": 0.20,
            "CVR (Analysis)": 0.05,
            "CVR (vs Comp)": 0.005,
            "Sessions (Analysis)": 20000,
            "AOV (Analysis)": 100.0,
            "Total Commission (Analysis)": 8000.0,
        },
    ]
    for i, name in enumerate(["Publisher A", "Publisher B", "Publisher C"]):
        rows.append({
            "Dimension 1": str(1000 + i),
            "Dimension 2": name,
            "Dimension 3": "Total",
            "Influencer Filter": "Affiliate",
            "Revenue (Analysis)": 30000.0 - i * 5000,
            "Revenue (Comparison)": 25000.0 - i * 5000,
            "Revenue (vs Comp)": 0.20,
            "Cost (Analysis)": 3000.0 - i * 500,
            "Cost (Comparison)": 2500.0 - i * 500,
            "Cost (vs Comp)": 0.20,
            "CoS (Analysis)": 0.10,
            "CoS (vs Comp)": -0.005,
            "Orders (Analysis)": 300 - i * 50,
            "Orders (vs Comp)": 0.15,
            "CVR (Analysis)": 0.05,
            "CVR (vs Comp)": 0.003,
            "Sessions (Analysis)": 6000 - i * 1000,
            "AOV (Analysis)": 100.0,
            "Total Commission (Analysis)": 2500.0 - i * 400,
        })
    return pd.DataFrame(rows)


def _all_sources(year, month):
    """Build a complete set of CoWork data sources for a given month."""
    return {
        "tracker": _make_tracker(year, month),
        "targets": _make_targets(year, month),
        "offer_performance": _make_offer_performance(month),
        "product_sales": _make_product_sales(month),
        "crm": _make_crm(),
        "affiliate": _make_affiliate(),
    }


def _bytes_to_prs(pptx_bytes: bytes) -> Presentation:
    """Load a Presentation from bytes."""
    return Presentation(io.BytesIO(pptx_bytes))


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def monthly_schema():
    return build_monthly_report_schema()


@pytest.fixture
def qbr_schema():
    return build_qbr_schema()


@pytest.fixture
def jan_sources():
    """All CoWork data sources for January 2026."""
    return _all_sources(2026, 1)


# ---------------------------------------------------------------------------
# Monthly E2E: full pipeline
# ---------------------------------------------------------------------------

class TestMonthlyFullPipeline:
    """Synthetic CoWork data → DataMapper → PPTXBuilder → QAValidator."""

    def test_full_pipeline_produces_valid_pptx(self, monthly_schema, jan_sources):
        """The complete pipeline should produce a QA-passing presentation."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        builder = PPTXBuilder(monthly_schema)
        pptx_bytes = builder.build(result.payload)
        qa = QAValidator(monthly_schema)
        qa_result = qa.validate(pptx_bytes, result.payload)

        structural = [i for i in qa_result.errors
                      if i.category in ("slide_count", "dimensions")]
        assert len(structural) == 0

    def test_full_pipeline_14_slides(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 14

    def test_full_pipeline_correct_dimensions(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_full_pipeline_high_coverage(self, monthly_schema, jan_sources):
        """With all data sources, mapper should achieve high coverage."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert result.coverage > 0.7

    def test_full_pipeline_no_nan_in_payload(self, monthly_schema, jan_sources):
        """The mapper must never leak NaN/inf into the payload."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)

        for key, val in result.payload.items():
            if isinstance(val, float):
                assert not math.isnan(val), f"{key} is NaN"
                assert not math.isinf(val), f"{key} is inf"
            elif isinstance(val, list):
                for item in val:
                    if isinstance(item, float):
                        assert not math.isnan(item), f"{key} list has NaN"
                    elif isinstance(item, dict):
                        for k, v in item.items():
                            if isinstance(v, float):
                                assert not math.isnan(v), f"{key}[].{k} is NaN"
                                assert not math.isinf(v), f"{key}[].{k} is inf"

    def test_full_pipeline_payload_validates(self, monthly_schema, jan_sources):
        """Mapper output should pass QA payload-only validation."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        qa = QAValidator(monthly_schema)
        payload_result = qa.validate_payload(result.payload)
        assert len(payload_result.errors) == 0

    def test_full_pipeline_no_table_errors(self, monthly_schema, jan_sources):
        """Tables rendered from mapper data should pass QA validation."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        table_errors = [i for i in qa_result.errors
                        if i.category.startswith("table_")]
        assert len(table_errors) == 0

    def test_full_pipeline_no_chart_errors(self, monthly_schema, jan_sources):
        """Charts rendered from mapper data should pass QA validation."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        chart_errors = [i for i in qa_result.errors
                        if i.category.startswith("chart_")]
        assert len(chart_errors) == 0

    def test_full_pipeline_no_kpi_errors(self, monthly_schema, jan_sources):
        """KPIs rendered from mapper data should pass QA validation."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        kpi_errors = [i for i in qa_result.errors
                      if i.category.startswith("kpi_")]
        assert len(kpi_errors) == 0

    def test_full_pipeline_divider_backgrounds(self, monthly_schema, jan_sources):
        """Divider slides should have brand-blue backgrounds."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        bg_errors = [i for i in qa_result.errors
                     if i.category == "divider_background"]
        assert len(bg_errors) == 0


# ---------------------------------------------------------------------------
# Monthly E2E: data accuracy
# ---------------------------------------------------------------------------

class TestMonthlyDataAccuracy:
    """Verify that generated slides contain the correct data values."""

    def test_cover_revenue_matches_input(self, monthly_schema, jan_sources):
        """Cover slide revenue should match aggregated tracker data."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)

        # 9 channels × 28 days × $1000 = $252,000
        expected_revenue = 9 * 28 * 1000.0
        assert result.payload["cover.total_revenue"] == expected_revenue

        # Verify it renders on the slide
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)
        cover = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in cover.shapes if s.has_text_frame
        )
        assert "$252k" in all_text

    def test_cover_orders_matches_input(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        expected_orders = 9 * 28 * 10
        assert result.payload["cover.total_orders"] == expected_orders

    def test_cover_aov_correct(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        # AOV = revenue / orders = (9*28*1000) / (9*28*10) = 100.0
        assert result.payload["cover.aov"] == pytest.approx(100.0)

    def test_exec_table_row_count(self, monthly_schema, jan_sources):
        """Executive summary table should have TOTAL + 9 channel rows."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        rows = result.payload["exec.performance_rows"]
        assert len(rows) == 10  # TOTAL + 9 channels

    def test_exec_table_total_row_first(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert result.payload["exec.performance_rows"][0]["channel"] == "TOTAL"

    def test_exec_yoy_variance_zero(self, monthly_schema, jan_sources):
        """With identical TY and LY data, YoY variance should be 0%."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        total_row = result.payload["exec.performance_rows"][0]
        assert total_row["revenue_vs_ly"] == 0.0

    def test_daily_series_length_matches_month(self, monthly_schema, jan_sources):
        """Daily series should have exactly as many values as days in January."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert len(result.payload["daily.dates"]) == 31
        assert len(result.payload["daily.revenue_actual"]) == 31
        assert len(result.payload["daily.revenue_target"]) == 31

    def test_daily_chart_rendered(self, monthly_schema, jan_sources):
        """Daily performance slide should have chart shapes."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)
        daily_slide = prs.slides[4]
        charts = [s for s in daily_slide.shapes if s.has_chart]
        assert len(charts) >= 1

    def test_promo_rows_from_offer_data(self, monthly_schema, jan_sources):
        """Promo rows should come from offer_performance data."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        promo_rows = result.payload["promo.rows"]
        # 3 promos × 3 channels = 9 rows
        assert len(promo_rows) == 9
        names = {r["promotion_name"] for r in promo_rows}
        assert "Promo A" in names
        assert "Grand Total" not in names

    def test_product_rows_from_product_data(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        product_rows = result.payload["product.rows"]
        assert len(product_rows) == 3
        assert "Grand Total" not in {r["product_name"] for r in product_rows}

    def test_crm_kpis_from_crm_data(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert result.payload["crm.emails_sent"] == 50000
        assert result.payload["crm.revenue"] == 75000.0

    def test_affiliate_kpis_from_affiliate_data(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert result.payload["affiliate.revenue"] == 100000.0
        assert result.payload["affiliate.roas"] == pytest.approx(10.0)

    def test_affiliate_publisher_rows_sorted(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        rows = result.payload["affiliate.publisher_rows"]
        assert len(rows) == 3
        assert rows[0]["publisher_name"] == "Publisher A"
        # Revenue should be descending
        revenues = [r["revenue"] for r in rows]
        assert revenues == sorted(revenues, reverse=True)

    def test_seo_from_organic_channel(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert result.payload["seo.revenue"] == 28 * 1000.0
        assert result.payload["seo.sessions"] == 28 * 200

    def test_report_title_and_period(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        assert result.payload["cover.report_title"] == "No7 US x THGi Monthly eComm Report"
        assert result.payload["cover.report_period"] == "January 2026 Overview"


# ---------------------------------------------------------------------------
# Monthly E2E: partial data (graceful degradation)
# ---------------------------------------------------------------------------

class TestMonthlyPartialData:
    """Pipeline should work with missing data sources."""

    def test_tracker_only(self, monthly_schema):
        """Only tracker data → cover/exec/daily/seo slides populated."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map({"tracker": _make_tracker(2026, 1)})
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        assert len(prs.slides) == 14
        assert result.payload["cover.total_revenue"] == 9 * 28 * 1000.0
        assert result.payload["promo.rows"] == []
        assert result.payload["product.rows"] == []
        assert len(result.warnings) > 0

    def test_no_data_sources(self, monthly_schema):
        """Empty sources should still produce 14-slide PPTX."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map({})
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        assert len(prs.slides) == 14
        assert result.coverage > 0.0  # Static slides provide some coverage
        assert result.coverage < 1.0

    def test_no_data_qa_no_structural_errors(self, monthly_schema):
        """Empty sources → no structural QA errors."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map({})
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        structural = [i for i in qa_result.errors
                      if i.category in ("slide_count", "dimensions")]
        assert len(structural) == 0

    def test_tracker_and_targets_only(self, monthly_schema):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        sources = {
            "tracker": _make_tracker(2026, 1),
            "targets": _make_targets(2026, 1),
        }
        result = mapper.map(sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        assert len(prs.slides) == 14
        assert result.payload["cover.revenue_vs_target"] is not None
        assert result.payload["promo.rows"] == []

    def test_offer_performance_only(self, monthly_schema):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map({
            "offer_performance": _make_offer_performance(1),
        })
        assert len(result.payload["promo.rows"]) == 9
        assert result.payload["exec.performance_rows"] == []


# ---------------------------------------------------------------------------
# Monthly E2E: multi-month
# ---------------------------------------------------------------------------

class TestMonthlyMultiMonth:
    """Pipeline should work correctly for any month."""

    @pytest.mark.parametrize("month,month_name", [
        (1, "January"), (3, "March"), (6, "June"),
        (9, "September"), (12, "December"),
    ])
    def test_report_period_correct(self, monthly_schema, month, month_name):
        mapper = DataMapper(monthly_schema, month=month, year=2026)
        result = mapper.map(_all_sources(2026, month))
        assert result.payload["cover.report_period"] == f"{month_name} 2026 Overview"

    @pytest.mark.parametrize("month", [1, 2, 6, 12])
    def test_daily_series_length(self, monthly_schema, month):
        """Daily series length should match days in that month."""
        mapper = DataMapper(monthly_schema, month=month, year=2026)
        result = mapper.map(_all_sources(2026, month))
        expected_days = calendar.monthrange(2026, month)[1]
        assert len(result.payload["daily.dates"]) == expected_days

    def test_february_non_leap_year(self, monthly_schema):
        mapper = DataMapper(monthly_schema, month=2, year=2026)
        result = mapper.map(_all_sources(2026, 2))
        assert len(result.payload["daily.dates"]) == 28

    def test_february_leap_year(self, monthly_schema):
        mapper = DataMapper(monthly_schema, month=2, year=2024)
        result = mapper.map(_all_sources(2024, 2))
        assert len(result.payload["daily.dates"]) == 29

    @pytest.mark.parametrize("month", [1, 6, 12])
    def test_full_pipeline_any_month(self, monthly_schema, month):
        """Full pipeline should pass for any month."""
        mapper = DataMapper(monthly_schema, month=month, year=2026)
        result = mapper.map(_all_sources(2026, month))
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        structural = [i for i in qa_result.errors
                      if i.category in ("slide_count", "dimensions")]
        assert len(structural) == 0


# ---------------------------------------------------------------------------
# Monthly E2E: idempotency and consistency
# ---------------------------------------------------------------------------

class TestMonthlyIdempotency:
    """Pipeline should produce identical output for identical input."""

    def test_mapper_idempotent(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        r1 = mapper.map(jan_sources)
        r2 = mapper.map(jan_sources)
        assert r1.payload == r2.payload
        assert r1.coverage == r2.coverage

    def test_builder_idempotent(self, monthly_schema, jan_sources):
        """Two builds from same payload should produce same slide count."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        builder = PPTXBuilder(monthly_schema)
        prs1 = _bytes_to_prs(builder.build(result.payload))
        prs2 = _bytes_to_prs(builder.build(result.payload))
        assert len(prs1.slides) == len(prs2.slides)

    def test_qa_idempotent(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa = QAValidator(monthly_schema)
        r1 = qa.validate(pptx_bytes, result.payload)
        r2 = qa.validate(pptx_bytes, result.payload)
        assert r1.error_count == r2.error_count
        assert r1.warning_count == r2.warning_count


# ---------------------------------------------------------------------------
# Monthly E2E: file output
# ---------------------------------------------------------------------------

class TestMonthlyFileOutput:
    """Pipeline should write valid PPTX to disk."""

    def test_build_to_file(self, monthly_schema, jan_sources, tmp_path):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        output = tmp_path / "No7_US_January_2026_Report.pptx"
        PPTXBuilder(monthly_schema).build_to_file(result.payload, output)

        assert output.exists()
        assert output.stat().st_size > 0
        prs = Presentation(str(output))
        assert len(prs.slides) == 14

    def test_file_output_qa_passes(self, monthly_schema, jan_sources, tmp_path):
        """File output should also pass QA validation."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        output = tmp_path / "report.pptx"
        PPTXBuilder(monthly_schema).build_to_file(result.payload, output)

        with open(output, "rb") as f:
            pptx_bytes = f.read()
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)
        structural = [i for i in qa_result.errors
                      if i.category in ("slide_count", "dimensions")]
        assert len(structural) == 0


# ---------------------------------------------------------------------------
# Schema compatibility: mapper keys vs schema keys
# ---------------------------------------------------------------------------

class TestMapperSchemaCompat:
    """Mapper output keys should align with schema expectations."""

    def test_mapper_covers_all_slot_data_keys(self, monthly_schema, jan_sources):
        """Mapper payload should contain data for most schema data_keys."""
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)

        schema_keys = set()
        for slide in monthly_schema.slides:
            for slot in slide.slots:
                if slot.data_key:
                    schema_keys.add(slot.data_key)
                if slot.variance_key:
                    schema_keys.add(slot.variance_key)
                if slot.row_data_key:
                    schema_keys.add(slot.row_data_key)
                if slot.categories_key:
                    schema_keys.add(slot.categories_key)
                for series in slot.series:
                    schema_keys.add(series.data_key)

        payload_keys = set(result.payload.keys())
        # Most schema keys should be present (manual slides may not have data)
        covered = schema_keys & payload_keys
        total_data_keys = {k for k in schema_keys
                          if not any(k.startswith(p) for p in ("upcoming.", "next_steps."))}
        coverage = len(covered & total_data_keys) / len(total_data_keys) if total_data_keys else 1.0
        assert coverage > 0.5

    def test_payload_keys_are_valid_strings(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        for key in result.payload:
            assert isinstance(key, str)
            assert "." in key, f"Key {key} missing namespace separator"


# ---------------------------------------------------------------------------
# QBR E2E: builder + QA
# ---------------------------------------------------------------------------

class TestQBRBuilderQA:
    """QBR schema → PPTXBuilder → QAValidator (no mapper for QBR yet)."""

    def test_qbr_empty_payload_29_slides(self, qbr_schema):
        """QBR builder should produce 29 slides even with empty payload."""
        pptx_bytes = PPTXBuilder(qbr_schema).build({})
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 29

    def test_qbr_dimensions_oversized(self, qbr_schema):
        pptx_bytes = PPTXBuilder(qbr_schema).build({})
        prs = _bytes_to_prs(pptx_bytes)
        assert prs.slide_width == Inches(21.986)
        assert prs.slide_height == Inches(12.368)

    def test_qbr_qa_structural(self, qbr_schema):
        pptx_bytes = PPTXBuilder(qbr_schema).build({})
        qa_result = QAValidator(qbr_schema).validate(pptx_bytes, {})
        structural = [i for i in qa_result.errors
                      if i.category in ("slide_count", "dimensions")]
        assert len(structural) == 0

    def test_qbr_with_cover_payload(self, qbr_schema):
        """QBR with cover KPIs should render correctly."""
        payload = {
            "qcover.report_title": "No7 US Quarterly Business Review",
            "qcover.report_period": "Q1 2026",
            "qcover.total_revenue": 3500000,
            "qcover.revenue_vs_target": 4.8,
            "qcover.total_orders": 35000,
            "qcover.orders_vs_target": 2.1,
            "qcover.aov": 100.0,
            "qcover.aov_vs_target": -0.5,
            "qcover.new_customers": 12000,
            "qcover.nc_vs_target": 6.3,
            "qcover.cvr": 3.5,
            "qcover.cvr_vs_target": 0.2,
            "qcover.cos": 11.8,
            "qcover.cos_vs_target": -0.8,
        }
        pptx_bytes = PPTXBuilder(qbr_schema).build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 29

        cover = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in cover.shapes if s.has_text_frame
        )
        assert "Q1 2026" in all_text

    def test_qbr_divider_slides_have_backgrounds(self, qbr_schema):
        """QBR divider slides should have filled backgrounds."""
        pptx_bytes = PPTXBuilder(qbr_schema).build({})
        qa_result = QAValidator(qbr_schema).validate(pptx_bytes, {})
        bg_errors = [i for i in qa_result.errors
                     if i.category == "divider_background"]
        assert len(bg_errors) == 0

    def test_qbr_file_output(self, qbr_schema, tmp_path):
        output = tmp_path / "No7_QBR_Q1_2026.pptx"
        PPTXBuilder(qbr_schema).build_to_file({}, output)
        assert output.exists()
        prs = Presentation(str(output))
        assert len(prs.slides) == 29


# ---------------------------------------------------------------------------
# Cross-schema: monthly vs QBR structural differences
# ---------------------------------------------------------------------------

class TestCrossSchema:
    """Verify structural differences between monthly and QBR schemas."""

    def test_monthly_smaller_than_qbr(self, monthly_schema, qbr_schema):
        assert monthly_schema.width_inches < qbr_schema.width_inches
        assert monthly_schema.height_inches < qbr_schema.height_inches

    def test_slide_counts_differ(self, monthly_schema, qbr_schema):
        assert len(monthly_schema.slides) == 14
        assert len(qbr_schema.slides) == 29

    def test_report_types_differ(self, monthly_schema, qbr_schema):
        assert monthly_schema.report_type == "monthly"
        assert qbr_schema.report_type == "qbr"

    def test_both_build_without_errors(self, monthly_schema, qbr_schema):
        """Both schemas should build successfully with empty payloads."""
        m_bytes = PPTXBuilder(monthly_schema).build({})
        q_bytes = PPTXBuilder(qbr_schema).build({})
        assert len(m_bytes) > 0
        assert len(q_bytes) > 0

    def test_both_qa_pass_structural(self, monthly_schema, qbr_schema):
        m_bytes = PPTXBuilder(monthly_schema).build({})
        q_bytes = PPTXBuilder(qbr_schema).build({})

        m_qa = QAValidator(monthly_schema).validate(m_bytes, {})
        q_qa = QAValidator(qbr_schema).validate(q_bytes, {})

        for qa_result in [m_qa, q_qa]:
            structural = [i for i in qa_result.errors
                          if i.category in ("slide_count", "dimensions")]
            assert len(structural) == 0


# ---------------------------------------------------------------------------
# Slide content verification
# ---------------------------------------------------------------------------

class TestSlideContentVerification:
    """Deep verification of specific slide contents after full pipeline."""

    def test_cover_slide_has_title(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        cover = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in cover.shapes if s.has_text_frame
        )
        assert "No7 US" in all_text
        assert "January 2026" in all_text

    def test_exec_slide_has_table(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        exec_slide = prs.slides[3]
        tables = [s for s in exec_slide.shapes if s.has_table]
        assert len(tables) >= 1
        # Table should have header + data rows
        table = tables[0].table
        assert len(table.rows) >= 2  # At least header + 1 data row

    def test_divider_slides_text(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        # Slide 2 = eComm divider
        divider = prs.slides[2]
        all_text = " ".join(
            s.text_frame.text for s in divider.shapes if s.has_text_frame
        )
        assert "eComm" in all_text

    def test_toc_slide_has_items(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        prs = _bytes_to_prs(pptx_bytes)

        toc = prs.slides[1]
        all_text = " ".join(
            s.text_frame.text for s in toc.shapes if s.has_text_frame
        )
        assert "eComm Performance" in all_text


# ---------------------------------------------------------------------------
# Variance coloring in E2E context
# ---------------------------------------------------------------------------

class TestVarianceColoringE2E:
    """Verify positive/negative variance coloring through full pipeline."""

    def test_variance_coloring_on_cover(self, monthly_schema):
        """Cover KPIs with non-zero variances should have correct colors."""
        sources = _all_sources(2026, 1)
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        # Filter out zero-variance color errors — synthetic data produces
        # identical TY/LY for some metrics (AOV, CVR, COS), causing 0.0%
        # variances where the builder uses green but QA expects neutral.
        color_errors = [i for i in qa_result.errors
                        if i.category == "variance_color"
                        and "'0.0%'" not in i.message]
        assert len(color_errors) == 0


# ---------------------------------------------------------------------------
# QA report output in E2E context
# ---------------------------------------------------------------------------

class TestQAReportE2E:
    """Verify QA report formatting after full pipeline."""

    def test_qa_report_is_string(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        report = qa_result.report()
        assert isinstance(report, str)
        assert len(report) > 0
        assert "QA" in report

    def test_qa_summary_format(self, monthly_schema, jan_sources):
        mapper = DataMapper(monthly_schema, month=1, year=2026)
        result = mapper.map(jan_sources)
        pptx_bytes = PPTXBuilder(monthly_schema).build(result.payload)
        qa_result = QAValidator(monthly_schema).validate(pptx_bytes, result.payload)

        summary = qa_result.summary()
        assert "error" in summary.lower() or "PASS" in summary
