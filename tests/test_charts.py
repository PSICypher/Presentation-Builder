"""Tests for the chart generation module."""

import math

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from src.generator.charts import (
    _build_category_chart_data,
    _build_doughnut_chart_data,
    _hex_to_rgb,
    _is_doughnut,
    _safe_value,
    add_chart,
    add_slide_charts,
)
from src.schema.models import (
    ChartSeries,
    ChartType,
    DataSlot,
    DesignSystem,
    Position,
    SlideSchema,
    SlideType,
    SlotType,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def slide():
    """Create a fresh presentation and return the first blank slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(slide_layout)


@pytest.fixture
def design():
    return DesignSystem()


# ---------------------------------------------------------------------------
# Slot builders
# ---------------------------------------------------------------------------

def _column_slot(
    name="test_chart",
    categories_key="test.categories",
    series=None,
    position=None,
):
    if series is None:
        series = [
            ChartSeries(name="Series A", data_key="test.series_a", color="#0065E0"),
            ChartSeries(name="Series B", data_key="test.series_b", color="#D1D5DB"),
        ]
    if position is None:
        position = Position(left=0.5, top=0.5, width=8.0, height=4.0)
    return DataSlot(
        name=name,
        slot_type=SlotType.CHART,
        data_key="test.chart",
        position=position,
        chart_type=ChartType.COLUMN_CLUSTERED,
        categories_key=categories_key,
        series=series,
    )


def _doughnut_slot(
    name="test_gauge",
    chart_type=ChartType.DOUGHNUT,
    series=None,
):
    if series is None:
        series = [
            ChartSeries(name="Achieved", data_key="test.achieved", color="#0065E0"),
            ChartSeries(name="Remaining", data_key="test.remaining", color="#D1D5DB"),
        ]
    return DataSlot(
        name=name,
        slot_type=SlotType.CHART,
        data_key="test.gauge",
        position=Position(left=0.5, top=5.0, width=2.0, height=1.5),
        chart_type=chart_type,
        series=series,
    )


def _line_slot(
    name="test_line",
    categories_key="test.categories",
    series=None,
):
    if series is None:
        series = [
            ChartSeries(name="Trend A", data_key="test.trend_a", color="#0065E0"),
            ChartSeries(name="Trend B", data_key="test.trend_b", color="#1C2B33"),
        ]
    return DataSlot(
        name=name,
        slot_type=SlotType.CHART,
        data_key="test.line",
        position=Position(left=0.5, top=0.5, width=8.0, height=4.0),
        chart_type=ChartType.LINE,
        categories_key=categories_key,
        series=series,
    )


# ---------------------------------------------------------------------------
# _hex_to_rgb
# ---------------------------------------------------------------------------

class TestHexToRgb:
    def test_black(self):
        rgb = _hex_to_rgb("#000000")
        assert (rgb[0], rgb[1], rgb[2]) == (0, 0, 0)

    def test_white(self):
        rgb = _hex_to_rgb("#FFFFFF")
        assert (rgb[0], rgb[1], rgb[2]) == (255, 255, 255)

    def test_brand_blue(self):
        rgb = _hex_to_rgb("#0065E0")
        assert (rgb[0], rgb[1], rgb[2]) == (0, 101, 224)

    def test_no_hash_prefix(self):
        rgb = _hex_to_rgb("0065E0")
        assert (rgb[0], rgb[1], rgb[2]) == (0, 101, 224)

    def test_lowercase(self):
        rgb = _hex_to_rgb("#ff0000")
        assert (rgb[0], rgb[1], rgb[2]) == (255, 0, 0)


# ---------------------------------------------------------------------------
# _safe_value
# ---------------------------------------------------------------------------

class TestSafeValue:
    def test_normal_float(self):
        assert _safe_value(42.5) == 42.5

    def test_normal_int(self):
        assert _safe_value(10) == 10.0

    def test_zero(self):
        assert _safe_value(0) == 0.0

    def test_negative(self):
        assert _safe_value(-5.0) == -5.0

    def test_none(self):
        assert _safe_value(None) == 0.0

    def test_nan(self):
        assert _safe_value(float("nan")) == 0.0

    def test_inf(self):
        assert _safe_value(float("inf")) == 0.0

    def test_neg_inf(self):
        assert _safe_value(float("-inf")) == 0.0

    def test_string(self):
        assert _safe_value("hello") == 0.0

    def test_list(self):
        assert _safe_value([1, 2]) == 0.0


# ---------------------------------------------------------------------------
# _is_doughnut
# ---------------------------------------------------------------------------

class TestIsDoughnut:
    def test_doughnut(self):
        assert _is_doughnut(ChartType.DOUGHNUT) is True

    def test_doughnut_exploded(self):
        assert _is_doughnut(ChartType.DOUGHNUT_EXPLODED) is True

    def test_column(self):
        assert _is_doughnut(ChartType.COLUMN_CLUSTERED) is False

    def test_line(self):
        assert _is_doughnut(ChartType.LINE) is False


# ---------------------------------------------------------------------------
# _build_category_chart_data
# ---------------------------------------------------------------------------

class TestBuildCategoryChartData:
    def test_normal(self):
        slot = _column_slot()
        payload = {
            "test.categories": ["Jan", "Feb", "Mar"],
            "test.series_a": [100, 200, 300],
            "test.series_b": [50, 100, 150],
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_missing_categories_key(self):
        slot = _column_slot(categories_key=None)
        data = _build_category_chart_data(slot, {})
        assert data is None

    def test_missing_categories_in_payload(self):
        slot = _column_slot()
        data = _build_category_chart_data(slot, {})
        assert data is None

    def test_empty_categories(self):
        slot = _column_slot()
        data = _build_category_chart_data(slot, {"test.categories": []})
        assert data is None

    def test_missing_series_data_uses_zeros(self):
        slot = _column_slot()
        payload = {"test.categories": ["Jan", "Feb"]}
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_short_series_padded(self):
        slot = _column_slot()
        payload = {
            "test.categories": ["Jan", "Feb", "Mar"],
            "test.series_a": [100],
            "test.series_b": [50, 100, 150],
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_long_series_truncated(self):
        slot = _column_slot()
        payload = {
            "test.categories": ["Jan", "Feb"],
            "test.series_a": [100, 200, 300, 400],
            "test.series_b": [50, 100],
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_no_series_definitions(self):
        slot = _column_slot(series=[])
        payload = {"test.categories": ["Jan", "Feb"]}
        data = _build_category_chart_data(slot, payload)
        assert data is None

    def test_none_values_in_series(self):
        slot = _column_slot()
        payload = {
            "test.categories": ["Jan", "Feb", "Mar"],
            "test.series_a": [100, None, 300],
            "test.series_b": [None, None, None],
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_nan_values_in_series(self):
        slot = _column_slot()
        payload = {
            "test.categories": ["A", "B"],
            "test.series_a": [float("nan"), 100],
            "test.series_b": [50, float("inf")],
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_single_category(self):
        slot = _column_slot()
        payload = {
            "test.categories": ["Only"],
            "test.series_a": [42],
            "test.series_b": [24],
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None

    def test_many_categories(self):
        slot = _column_slot()
        cats = [f"Day {i}" for i in range(1, 32)]
        payload = {
            "test.categories": cats,
            "test.series_a": list(range(31)),
            "test.series_b": list(range(31)),
        }
        data = _build_category_chart_data(slot, payload)
        assert data is not None


# ---------------------------------------------------------------------------
# _build_doughnut_chart_data
# ---------------------------------------------------------------------------

class TestBuildDoughnutChartData:
    def test_normal(self):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0.75, "test.remaining": 0.25}
        data = _build_doughnut_chart_data(slot, payload)
        assert data is not None

    def test_all_zeros(self):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0, "test.remaining": 0}
        data = _build_doughnut_chart_data(slot, payload)
        assert data is None

    def test_all_missing(self):
        slot = _doughnut_slot()
        data = _build_doughnut_chart_data(slot, {})
        assert data is None

    def test_partial_data(self):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0.5}
        data = _build_doughnut_chart_data(slot, payload)
        assert data is not None

    def test_no_series(self):
        slot = _doughnut_slot(series=[])
        data = _build_doughnut_chart_data(slot, {})
        assert data is None

    def test_none_values(self):
        slot = _doughnut_slot()
        payload = {"test.achieved": None, "test.remaining": None}
        data = _build_doughnut_chart_data(slot, payload)
        assert data is None

    def test_single_slice(self):
        slot = _doughnut_slot(series=[
            ChartSeries(name="Full", data_key="test.full", color="#0065E0"),
        ])
        data = _build_doughnut_chart_data(slot, {"test.full": 1.0})
        assert data is not None

    def test_many_slices(self):
        series = [
            ChartSeries(name=f"Slice {i}", data_key=f"test.s{i}")
            for i in range(5)
        ]
        slot = _doughnut_slot(series=series)
        payload = {f"test.s{i}": 20.0 for i in range(5)}
        data = _build_doughnut_chart_data(slot, payload)
        assert data is not None


# ---------------------------------------------------------------------------
# add_chart — column charts
# ---------------------------------------------------------------------------

class TestAddChartColumn:
    def test_basic_column(self, slide, design):
        slot = _column_slot()
        payload = {
            "test.categories": ["Jan", "Feb", "Mar"],
            "test.series_a": [100, 200, 300],
            "test.series_b": [50, 100, 150],
        }
        result = add_chart(slide, slot, payload, design)
        assert result is True
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 1

    def test_column_returns_false_no_data(self, slide, design):
        slot = _column_slot()
        result = add_chart(slide, slot, {}, design)
        assert result is False
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 0

    def test_column_position(self, slide, design):
        pos = Position(left=1.0, top=2.0, width=6.0, height=3.0)
        slot = _column_slot(position=pos)
        payload = {
            "test.categories": ["A"],
            "test.series_a": [1],
            "test.series_b": [2],
        }
        add_chart(slide, slot, payload, design)
        shape = [s for s in slide.shapes if s.has_chart][0]
        assert shape.left == Inches(1.0)
        assert shape.top == Inches(2.0)
        assert shape.width == Inches(6.0)
        assert shape.height == Inches(3.0)

    def test_column_series_colors(self, slide, design):
        slot = _column_slot()
        payload = {
            "test.categories": ["A"],
            "test.series_a": [100],
            "test.series_b": [50],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        plot = chart.plots[0]
        assert plot.series[0].format.fill.fore_color.rgb == _hex_to_rgb("#0065E0")
        assert plot.series[1].format.fill.fore_color.rgb == _hex_to_rgb("#D1D5DB")

    def test_column_series_count(self, slide, design):
        slot = _column_slot()
        payload = {
            "test.categories": ["A", "B"],
            "test.series_a": [1, 2],
            "test.series_b": [3, 4],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert len(chart.plots[0].series) == 2

    def test_column_legend_multi_series(self, slide, design):
        slot = _column_slot()
        payload = {
            "test.categories": ["A"],
            "test.series_a": [100],
            "test.series_b": [50],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.has_legend is True

    def test_column_no_legend_single_series(self, slide, design):
        slot = _column_slot(series=[
            ChartSeries(name="Only", data_key="test.only", color="#0065E0"),
        ])
        payload = {"test.categories": ["A", "B"], "test.only": [1, 2]}
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.has_legend is False

    def test_column_font_styling(self, slide, design):
        slot = _column_slot()
        payload = {
            "test.categories": ["A"],
            "test.series_a": [1],
            "test.series_b": [2],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.font.name == "DM Sans"
        assert chart.font.size == Pt(9.0)

    def test_column_three_series(self, slide, design):
        slot = _column_slot(series=[
            ChartSeries(name="A", data_key="test.a", color="#0065E0"),
            ChartSeries(name="B", data_key="test.b", color="#D1D5DB"),
            ChartSeries(name="C", data_key="test.c", color="#1C2B33"),
        ])
        payload = {
            "test.categories": ["X", "Y"],
            "test.a": [10, 20],
            "test.b": [30, 40],
            "test.c": [50, 60],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert len(chart.plots[0].series) == 3

    def test_column_no_series_color(self, slide, design):
        slot = _column_slot(series=[
            ChartSeries(name="NoColor", data_key="test.nc"),
        ])
        payload = {"test.categories": ["A"], "test.nc": [10]}
        result = add_chart(slide, slot, payload, design)
        assert result is True


# ---------------------------------------------------------------------------
# add_chart — doughnut charts
# ---------------------------------------------------------------------------

class TestAddChartDoughnut:
    def test_basic_doughnut(self, slide, design):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0.75, "test.remaining": 0.25}
        result = add_chart(slide, slot, payload, design)
        assert result is True
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 1

    def test_doughnut_exploded(self, slide, design):
        slot = _doughnut_slot(chart_type=ChartType.DOUGHNUT_EXPLODED)
        payload = {"test.achieved": 0.6, "test.remaining": 0.4}
        result = add_chart(slide, slot, payload, design)
        assert result is True

    def test_doughnut_no_legend(self, slide, design):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0.75, "test.remaining": 0.25}
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.has_legend is False

    def test_doughnut_point_colors(self, slide, design):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0.75, "test.remaining": 0.25}
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        series = chart.plots[0].series[0]
        assert series.points[0].format.fill.fore_color.rgb == _hex_to_rgb("#0065E0")
        assert series.points[1].format.fill.fore_color.rgb == _hex_to_rgb("#D1D5DB")

    def test_doughnut_missing_data(self, slide, design):
        slot = _doughnut_slot()
        result = add_chart(slide, slot, {}, design)
        assert result is False

    def test_doughnut_all_zeros(self, slide, design):
        slot = _doughnut_slot()
        payload = {"test.achieved": 0, "test.remaining": 0}
        result = add_chart(slide, slot, payload, design)
        assert result is False


# ---------------------------------------------------------------------------
# add_chart — line charts
# ---------------------------------------------------------------------------

class TestAddChartLine:
    def test_basic_line(self, slide, design):
        slot = _line_slot()
        payload = {
            "test.categories": ["Jan", "Feb", "Mar", "Apr"],
            "test.trend_a": [10, 20, 30, 40],
            "test.trend_b": [5, 15, 25, 35],
        }
        result = add_chart(slide, slot, payload, design)
        assert result is True
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 1

    def test_line_series_line_color(self, slide, design):
        slot = _line_slot()
        payload = {
            "test.categories": ["A"],
            "test.trend_a": [10],
            "test.trend_b": [5],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        plot = chart.plots[0]
        assert plot.series[0].format.line.color.rgb == _hex_to_rgb("#0065E0")
        assert plot.series[1].format.line.color.rgb == _hex_to_rgb("#1C2B33")

    def test_line_legend(self, slide, design):
        slot = _line_slot()
        payload = {
            "test.categories": ["A"],
            "test.trend_a": [10],
            "test.trend_b": [5],
        }
        add_chart(slide, slot, payload, design)
        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.has_legend is True

    def test_line_missing_data(self, slide, design):
        slot = _line_slot()
        result = add_chart(slide, slot, {}, design)
        assert result is False


# ---------------------------------------------------------------------------
# add_chart — error handling
# ---------------------------------------------------------------------------

class TestAddChartErrors:
    def test_wrong_slot_type(self, slide, design):
        slot = DataSlot(
            name="text_slot",
            slot_type=SlotType.TEXT,
            data_key="test.text",
            position=Position(left=0, top=0, width=1, height=1),
        )
        with pytest.raises(ValueError, match="not a CHART type"):
            add_chart(slide, slot, {}, design)

    def test_missing_chart_type(self, slide, design):
        slot = DataSlot(
            name="bad_chart",
            slot_type=SlotType.CHART,
            data_key="test.chart",
            position=Position(left=0, top=0, width=1, height=1),
            chart_type=None,
        )
        with pytest.raises(ValueError, match="no chart_type"):
            add_chart(slide, slot, {}, design)

    def test_kpi_slot_raises(self, slide, design):
        slot = DataSlot(
            name="kpi",
            slot_type=SlotType.KPI_VALUE,
            data_key="test.kpi",
            position=Position(left=0, top=0, width=1, height=1),
        )
        with pytest.raises(ValueError, match="not a CHART type"):
            add_chart(slide, slot, {}, design)

    def test_table_slot_raises(self, slide, design):
        slot = DataSlot(
            name="table",
            slot_type=SlotType.TABLE,
            data_key="test.table",
            position=Position(left=0, top=0, width=1, height=1),
        )
        with pytest.raises(ValueError, match="not a CHART type"):
            add_chart(slide, slot, {}, design)


# ---------------------------------------------------------------------------
# Monthly report integration — daily performance charts
# ---------------------------------------------------------------------------

class TestMonthlyReportCharts:
    def test_daily_column_chart(self, slide, design):
        slot = DataSlot(
            name="daily_chart",
            slot_type=SlotType.CHART,
            data_key="daily.chart",
            position=Position(left=0.3, top=0.9, width=8.5, height=4.5),
            chart_type=ChartType.COLUMN_CLUSTERED,
            categories_key="daily.dates",
            series=[
                ChartSeries(name="Revenue", data_key="daily.revenue_actual",
                            color="#0065E0"),
                ChartSeries(name="Target", data_key="daily.revenue_target",
                            color="#D1D5DB"),
                ChartSeries(name="LY", data_key="daily.revenue_ly",
                            color="#1C2B33"),
            ],
        )
        payload = {
            "daily.dates": [f"1/{d}" for d in range(1, 32)],
            "daily.revenue_actual": [9000.0] * 28 + [0.0] * 3,
            "daily.revenue_target": [8500.0] * 31,
            "daily.revenue_ly": [7500.0] * 28 + [0.0] * 3,
        }
        result = add_chart(slide, slot, payload, design)
        assert result is True

        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.has_legend is True
        assert len(chart.plots[0].series) == 3

    def test_revenue_gauge_doughnut(self, slide, design):
        slot = DataSlot(
            name="revenue_gauge",
            slot_type=SlotType.CHART,
            data_key="daily.revenue_gauge",
            position=Position(left=0.5, top=5.5, width=2.0, height=1.5),
            chart_type=ChartType.DOUGHNUT,
            series=[
                ChartSeries(name="Achieved",
                            data_key="daily.revenue_achieved_pct",
                            color="#0065E0"),
                ChartSeries(name="Remaining",
                            data_key="daily.revenue_remaining_pct",
                            color="#D1D5DB"),
            ],
        )
        payload = {
            "daily.revenue_achieved_pct": 0.85,
            "daily.revenue_remaining_pct": 0.15,
        }
        result = add_chart(slide, slot, payload, design)
        assert result is True

        chart = [s for s in slide.shapes if s.has_chart][0].chart
        assert chart.has_legend is False
        assert len(list(chart.plots[0].series[0].points)) == 2


# ---------------------------------------------------------------------------
# add_slide_charts
# ---------------------------------------------------------------------------

class TestAddSlideCharts:
    def test_slide_with_mixed_slots(self, slide, design):
        schema = SlideSchema(
            index=4,
            name="daily_performance",
            title="Daily Performance",
            slide_type=SlideType.DATA,
            data_source="tracker",
            slots=[
                DataSlot(
                    name="slide_title",
                    slot_type=SlotType.TEXT,
                    data_key="daily.title",
                    position=Position(left=0.3, top=0.2, width=12.0, height=0.5),
                ),
                _column_slot(name="daily_chart"),
                _doughnut_slot(name="revenue_gauge"),
            ],
        )
        payload = {
            "test.categories": ["A", "B"],
            "test.series_a": [1, 2],
            "test.series_b": [3, 4],
            "test.achieved": 0.7,
            "test.remaining": 0.3,
        }
        added = add_slide_charts(slide, schema, payload, design)
        assert "daily_chart" in added
        assert "revenue_gauge" in added
        assert "slide_title" not in added
        assert len(added) == 2

    def test_slide_with_no_charts(self, slide, design):
        schema = SlideSchema(
            index=0,
            name="cover",
            title="Cover",
            slide_type=SlideType.COVER,
            data_source="tracker",
            slots=[
                DataSlot(
                    name="title",
                    slot_type=SlotType.TEXT,
                    data_key="cover.title",
                    position=Position(left=0.5, top=0.4, width=12.0, height=0.8),
                ),
            ],
        )
        added = add_slide_charts(slide, schema, {}, design)
        assert added == []

    def test_slide_with_missing_chart_data(self, slide, design):
        schema = SlideSchema(
            index=4,
            name="daily",
            title="Daily",
            slide_type=SlideType.DATA,
            data_source="tracker",
            slots=[_column_slot(name="chart1")],
        )
        added = add_slide_charts(slide, schema, {}, design)
        assert added == []

    def test_slide_returns_only_added(self, slide, design):
        schema = SlideSchema(
            index=4,
            name="daily",
            title="Daily",
            slide_type=SlideType.DATA,
            data_source="tracker",
            slots=[
                _column_slot(name="has_data"),
                _column_slot(name="no_data", categories_key="missing.cats"),
            ],
        )
        payload = {
            "test.categories": ["A"],
            "test.series_a": [1],
            "test.series_b": [2],
        }
        added = add_slide_charts(slide, schema, payload, design)
        assert added == ["has_data"]


# ---------------------------------------------------------------------------
# Multiple charts on one slide
# ---------------------------------------------------------------------------

class TestMultipleCharts:
    def test_two_charts_on_slide(self, slide, design):
        col_slot = _column_slot()
        dnut_slot = _doughnut_slot()
        payload = {
            "test.categories": ["A", "B"],
            "test.series_a": [1, 2],
            "test.series_b": [3, 4],
            "test.achieved": 0.7,
            "test.remaining": 0.3,
        }
        r1 = add_chart(slide, col_slot, payload, design)
        r2 = add_chart(slide, dnut_slot, payload, design)
        assert r1 is True
        assert r2 is True
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 2

    def test_three_charts_on_slide(self, slide, design):
        payload = {
            "test.categories": ["A", "B"],
            "test.series_a": [1, 2],
            "test.series_b": [3, 4],
            "test.trend_a": [10, 20],
            "test.trend_b": [5, 15],
            "test.achieved": 0.8,
            "test.remaining": 0.2,
        }
        add_chart(slide, _column_slot(), payload, design)
        add_chart(slide, _line_slot(), payload, design)
        add_chart(slide, _doughnut_slot(), payload, design)
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 3


# ---------------------------------------------------------------------------
# PPTX file round-trip
# ---------------------------------------------------------------------------

class TestPptxRoundTrip:
    def test_save_and_reload(self, slide, design, tmp_path):
        slot = _column_slot()
        payload = {
            "test.categories": ["Q1", "Q2", "Q3", "Q4"],
            "test.series_a": [100, 200, 150, 250],
            "test.series_b": [80, 180, 120, 220],
        }
        add_chart(slide, slot, payload, design)

        pptx_path = tmp_path / "test_chart.pptx"
        slide.part.package.save(str(pptx_path))

        reloaded = Presentation(str(pptx_path))
        reloaded_slide = reloaded.slides[0]
        chart_shapes = [s for s in reloaded_slide.shapes if s.has_chart]
        assert len(chart_shapes) == 1
        assert len(chart_shapes[0].chart.plots[0].series) == 2
