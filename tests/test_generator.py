"""Tests for the PPTX builder engine."""

import io
import math

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from src.generator.pptx_builder import (
    PPTXBuilder,
    build_presentation,
    _hex_to_rgb,
    _is_missing,
    _format_slot_value,
)
from src.schema.models import (
    ChartSeries,
    ChartType,
    DataSlot,
    DesignSystem,
    FontSpec,
    FormatRule,
    FormatType,
    Position,
    SlideSchema,
    SlideType,
    SlotType,
    TableColumn,
    TemplateSchema,
)
from src.schema.monthly_report import build_monthly_report_schema


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def design():
    return DesignSystem()


@pytest.fixture
def minimal_schema(design):
    """A minimal 1-slide schema for focused testing."""
    return TemplateSchema(
        name="Test Report",
        report_type="monthly",
        width_inches=13.333,
        height_inches=7.5,
        design=design,
        slides=[
            SlideSchema(
                index=0,
                name="test_slide",
                title="Test Slide",
                slide_type=SlideType.DATA,
                data_source="test",
                slots=[],
            ),
        ],
    )


@pytest.fixture
def full_schema():
    """The full 14-slide monthly report schema."""
    return build_monthly_report_schema()


def _bytes_to_prs(pptx_bytes: bytes) -> Presentation:
    """Load a Presentation from bytes."""
    return Presentation(io.BytesIO(pptx_bytes))


# ---------------------------------------------------------------------------
# Helper tests
# ---------------------------------------------------------------------------

class TestHexToRgb:
    def test_black(self):
        assert _hex_to_rgb("#000000") == RGBColor(0, 0, 0)

    def test_white(self):
        assert _hex_to_rgb("#FFFFFF") == RGBColor(255, 255, 255)

    def test_brand_blue(self):
        assert _hex_to_rgb("#0065E0") == RGBColor(0, 101, 224)

    def test_no_hash(self):
        assert _hex_to_rgb("FF0000") == RGBColor(255, 0, 0)


class TestIsMissing:
    def test_none(self):
        assert _is_missing(None) is True

    def test_nan(self):
        assert _is_missing(float("nan")) is True

    def test_zero(self):
        assert _is_missing(0) is False

    def test_empty_string(self):
        assert _is_missing("") is False

    def test_number(self):
        assert _is_missing(42.5) is False

    def test_list(self):
        assert _is_missing([]) is False


class TestFormatSlotValue:
    def test_currency(self):
        rule = FormatRule(FormatType.CURRENCY)
        assert _format_slot_value(1234567, rule) == "$1.2m"

    def test_percentage(self):
        rule = FormatRule(FormatType.PERCENTAGE)
        assert _format_slot_value(3.6, rule) == "3.6%"

    def test_variance_percentage(self):
        rule = FormatRule(FormatType.VARIANCE_PERCENTAGE)
        assert _format_slot_value(5.2, rule) == "+5.2%"

    def test_integer(self):
        rule = FormatRule(FormatType.INTEGER)
        assert _format_slot_value(1234, rule) == "1,234"

    def test_none_returns_na(self):
        rule = FormatRule(FormatType.CURRENCY)
        assert _format_slot_value(None, rule) == "N/A"

    def test_nan_returns_na(self):
        rule = FormatRule(FormatType.CURRENCY)
        assert _format_slot_value(float("nan"), rule) == "N/A"

    def test_no_format_rule(self):
        assert _format_slot_value(42, None) == "42"

    def test_no_format_rule_none(self):
        assert _format_slot_value(None, None) == "N/A"


# ---------------------------------------------------------------------------
# Slide dimension tests
# ---------------------------------------------------------------------------

class TestSlideDimensions:
    def test_standard_16_9(self, minimal_schema):
        builder = PPTXBuilder(minimal_schema)
        pptx_bytes = builder.build({})
        prs = _bytes_to_prs(pptx_bytes)
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_slide_count_matches_schema(self, minimal_schema):
        builder = PPTXBuilder(minimal_schema)
        pptx_bytes = builder.build({})
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 1

    def test_full_schema_slide_count(self, full_schema):
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build({})
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 14


# ---------------------------------------------------------------------------
# KPI rendering tests
# ---------------------------------------------------------------------------

class TestKPIRendering:
    def _kpi_schema(self, design):
        return TemplateSchema(
            name="KPI Test",
            report_type="monthly",
            width_inches=13.333,
            height_inches=7.5,
            design=design,
            slides=[
                SlideSchema(
                    index=0,
                    name="kpi_slide",
                    title="KPI Slide",
                    slide_type=SlideType.DATA,
                    data_source="test",
                    slots=[
                        DataSlot(
                            name="revenue",
                            slot_type=SlotType.KPI_VALUE,
                            data_key="test.revenue",
                            position=Position(left=0.5, top=1.0, width=2.0, height=1.5),
                            font=FontSpec(name="DM Sans", size_pt=48.0, bold=True),
                            format_rule=FormatRule(FormatType.CURRENCY),
                            label="Revenue",
                            variance_key="test.revenue_var",
                        ),
                    ],
                ),
            ],
        )

    def test_kpi_renders_with_data(self, design):
        schema = self._kpi_schema(design)
        payload = {"test.revenue": 209200, "test.revenue_var": 5.2}
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        # Should have a textbox shape
        assert len(slide.shapes) >= 1
        # Check text content includes formatted revenue and variance
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "$209.2k" in all_text
        assert "+5.2%" in all_text
        assert "Revenue" in all_text

    def test_kpi_renders_missing_data(self, design):
        schema = self._kpi_schema(design)
        payload = {}  # No data
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "N/A" in all_text
        assert "Revenue" in all_text

    def test_kpi_negative_variance_color(self, design):
        schema = self._kpi_schema(design)
        payload = {"test.revenue": 100000, "test.revenue_var": -3.1}
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        # Find the variance text
        found_negative = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "-3.1%" in run.text:
                            assert run.font.color.rgb == RGBColor(0xCC, 0x00, 0x00)
                            found_negative = True
        assert found_negative, "Negative variance text not found"


# ---------------------------------------------------------------------------
# Table rendering tests
# ---------------------------------------------------------------------------

class TestTableRendering:
    def _table_schema(self, design):
        return TemplateSchema(
            name="Table Test",
            report_type="monthly",
            width_inches=13.333,
            height_inches=7.5,
            design=design,
            slides=[
                SlideSchema(
                    index=0,
                    name="table_slide",
                    title="Table Slide",
                    slide_type=SlideType.DATA,
                    data_source="test",
                    slots=[
                        DataSlot(
                            name="test_table",
                            slot_type=SlotType.TABLE,
                            data_key="test.table",
                            position=Position(left=0.3, top=0.9, width=12.0, height=4.0),
                            row_data_key="test.rows",
                            columns=[
                                TableColumn(
                                    header="Channel",
                                    data_key="channel",
                                    width_inches=2.0,
                                    alignment="left",
                                ),
                                TableColumn(
                                    header="Revenue",
                                    data_key="revenue",
                                    width_inches=1.5,
                                    format_rule=FormatRule(FormatType.CURRENCY),
                                    alignment="right",
                                ),
                                TableColumn(
                                    header="vs Target",
                                    data_key="vs_target",
                                    width_inches=1.0,
                                    format_rule=FormatRule(FormatType.VARIANCE_PERCENTAGE),
                                    alignment="right",
                                ),
                            ],
                        ),
                    ],
                ),
            ],
        )

    def test_table_renders_with_data(self, design):
        schema = self._table_schema(design)
        payload = {
            "test.rows": [
                {"channel": "DIRECT", "revenue": 45000, "vs_target": 3.2},
                {"channel": "PPC", "revenue": 32000, "vs_target": -1.5},
            ],
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]

        # Find the table shape
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1

        table = tables[0].table
        assert len(table.rows) == 3  # 1 header + 2 data
        assert len(table.columns) == 3
        assert table.cell(0, 0).text == "Channel"
        assert table.cell(0, 1).text == "Revenue"
        assert table.cell(1, 0).text == "DIRECT"
        assert table.cell(1, 1).text == "$45k"
        assert table.cell(2, 0).text == "PPC"

    def test_table_renders_empty_data(self, design):
        schema = self._table_schema(design)
        payload = {}  # No row data
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        # Should not crash — falls back to text placeholder
        assert len(prs.slides) == 1

    def test_table_variance_coloring(self, design):
        schema = self._table_schema(design)
        payload = {
            "test.rows": [
                {"channel": "DIRECT", "revenue": 50000, "vs_target": 5.0},
                {"channel": "PPC", "revenue": 30000, "vs_target": -2.5},
            ],
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        table = [s for s in prs.slides[0].shapes if s.has_table][0].table

        # Check positive variance cell color
        pos_cell = table.cell(1, 2)
        for p in pos_cell.text_frame.paragraphs:
            for run in p.runs:
                if "+5.0%" in run.text:
                    assert run.font.color.rgb == RGBColor(0x00, 0xAA, 0x00)

        # Check negative variance cell color
        neg_cell = table.cell(2, 2)
        for p in neg_cell.text_frame.paragraphs:
            for run in p.runs:
                if "-2.5%" in run.text:
                    assert run.font.color.rgb == RGBColor(0xCC, 0x00, 0x00)

    def test_table_column_widths(self, design):
        schema = self._table_schema(design)
        payload = {
            "test.rows": [{"channel": "X", "revenue": 100, "vs_target": 0}],
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        table = [s for s in prs.slides[0].shapes if s.has_table][0].table

        assert table.columns[0].width == Inches(2.0)
        assert table.columns[1].width == Inches(1.5)
        assert table.columns[2].width == Inches(1.0)


# ---------------------------------------------------------------------------
# Chart rendering tests
# ---------------------------------------------------------------------------

class TestChartRendering:
    def _column_chart_schema(self, design):
        return TemplateSchema(
            name="Chart Test",
            report_type="monthly",
            width_inches=13.333,
            height_inches=7.5,
            design=design,
            slides=[
                SlideSchema(
                    index=0,
                    name="chart_slide",
                    title="Chart Slide",
                    slide_type=SlideType.DATA,
                    data_source="test",
                    slots=[
                        DataSlot(
                            name="daily_chart",
                            slot_type=SlotType.CHART,
                            data_key="test.chart",
                            position=Position(left=0.3, top=0.9, width=8.0, height=4.0),
                            chart_type=ChartType.COLUMN_CLUSTERED,
                            categories_key="test.dates",
                            series=[
                                ChartSeries(
                                    name="Revenue",
                                    data_key="test.revenue_series",
                                    color="#0065E0",
                                ),
                                ChartSeries(
                                    name="Target",
                                    data_key="test.target_series",
                                    color="#D1D5DB",
                                ),
                            ],
                        ),
                    ],
                ),
            ],
        )

    def _doughnut_chart_schema(self, design):
        return TemplateSchema(
            name="Doughnut Test",
            report_type="monthly",
            width_inches=13.333,
            height_inches=7.5,
            design=design,
            slides=[
                SlideSchema(
                    index=0,
                    name="gauge_slide",
                    title="Gauge Slide",
                    slide_type=SlideType.DATA,
                    data_source="test",
                    slots=[
                        DataSlot(
                            name="gauge",
                            slot_type=SlotType.CHART,
                            data_key="test.gauge",
                            position=Position(left=0.5, top=5.5, width=2.0, height=1.5),
                            chart_type=ChartType.DOUGHNUT,
                            series=[
                                ChartSeries(
                                    name="Achieved",
                                    data_key="test.achieved",
                                    color="#0065E0",
                                ),
                                ChartSeries(
                                    name="Remaining",
                                    data_key="test.remaining",
                                    color="#D1D5DB",
                                ),
                            ],
                        ),
                    ],
                ),
            ],
        )

    def test_column_chart_renders(self, design):
        schema = self._column_chart_schema(design)
        payload = {
            "test.dates": ["1/1", "1/2", "1/3"],
            "test.revenue_series": [12000, 15000, 18000],
            "test.target_series": [15000, 15000, 15000],
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]

        charts = [s for s in slide.shapes if s.has_chart]
        assert len(charts) == 1
        chart = charts[0].chart
        assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    def test_column_chart_series_count(self, design):
        schema = self._column_chart_schema(design)
        payload = {
            "test.dates": ["1/1", "1/2"],
            "test.revenue_series": [10000, 20000],
            "test.target_series": [15000, 15000],
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        chart = [s for s in prs.slides[0].shapes if s.has_chart][0].chart
        assert len(chart.series) == 2

    def test_doughnut_chart_renders(self, design):
        schema = self._doughnut_chart_schema(design)
        payload = {
            "test.achieved": 75.0,
            "test.remaining": 25.0,
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        charts = [s for s in prs.slides[0].shapes if s.has_chart]
        assert len(charts) == 1
        chart = charts[0].chart
        assert chart.chart_type == XL_CHART_TYPE.DOUGHNUT

    def test_chart_missing_data(self, design):
        schema = self._column_chart_schema(design)
        payload = {}  # No chart data
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        # Should render with zeros, not crash
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Text rendering tests
# ---------------------------------------------------------------------------

class TestTextRendering:
    def _text_schema(self, design):
        return TemplateSchema(
            name="Text Test",
            report_type="monthly",
            width_inches=13.333,
            height_inches=7.5,
            design=design,
            slides=[
                SlideSchema(
                    index=0,
                    name="text_slide",
                    title="Text Slide",
                    slide_type=SlideType.DATA,
                    data_source="test",
                    slots=[
                        DataSlot(
                            name="title",
                            slot_type=SlotType.TEXT,
                            data_key="test.title",
                            position=Position(left=0.3, top=0.2, width=12.0, height=0.5),
                            font=FontSpec(name="DM Sans", size_pt=24.0, bold=True),
                        ),
                        DataSlot(
                            name="body",
                            slot_type=SlotType.TEXT,
                            data_key="test.body",
                            position=Position(left=0.3, top=1.0, width=12.0, height=5.0),
                            font=FontSpec(name="DM Sans", size_pt=14.0),
                        ),
                    ],
                ),
            ],
        )

    def test_text_renders(self, design):
        schema = self._text_schema(design)
        payload = {
            "test.title": "Executive Summary",
            "test.body": "Revenue increased by 5%.",
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        assert "Executive Summary" in all_text
        assert "Revenue increased by 5%" in all_text

    def test_text_list_rendering(self, design):
        schema = self._text_schema(design)
        payload = {
            "test.title": "TOC",
            "test.body": ["Item 1", "Item 2", "Item 3"],
        }
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        assert "Item 1" in all_text
        assert "Item 2" in all_text
        assert "Item 3" in all_text

    def test_text_missing_data(self, design):
        schema = self._text_schema(design)
        payload = {}
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Section divider tests
# ---------------------------------------------------------------------------

class TestSectionDivider:
    def _divider_schema(self, design):
        return TemplateSchema(
            name="Divider Test",
            report_type="monthly",
            width_inches=13.333,
            height_inches=7.5,
            design=design,
            slides=[
                SlideSchema(
                    index=0,
                    name="divider",
                    title="eComm Performance",
                    slide_type=SlideType.SECTION_DIVIDER,
                    data_source="static",
                    is_static=True,
                    slots=[
                        DataSlot(
                            name="section_title",
                            slot_type=SlotType.SECTION_DIVIDER,
                            data_key="divider.title",
                            position=Position(left=0.0, top=0.0, width=13.333, height=7.5),
                            font=FontSpec(
                                name="DM Sans", size_pt=36.0,
                                bold=True, color="#FFFFFF",
                            ),
                        ),
                    ],
                ),
            ],
        )

    def test_divider_renders_with_text(self, design):
        schema = self._divider_schema(design)
        payload = {"divider.title": "eComm Performance"}
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        assert "eComm Performance" in all_text

    def test_divider_has_background_fill(self, design):
        schema = self._divider_schema(design)
        payload = {"divider.title": "Test"}
        builder = PPTXBuilder(schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)
        slide = prs.slides[0]
        # Background should have solid fill with brand blue
        bg = slide.background
        fill = bg.fill
        assert fill.fore_color.rgb == RGBColor(0x00, 0x65, 0xE0)


# ---------------------------------------------------------------------------
# Convenience function tests
# ---------------------------------------------------------------------------

class TestBuildPresentation:
    def test_build_presentation_returns_bytes(self, minimal_schema):
        result = build_presentation(minimal_schema, {})
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_build_presentation_is_valid_pptx(self, minimal_schema):
        result = build_presentation(minimal_schema, {})
        prs = _bytes_to_prs(result)
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Build to file tests
# ---------------------------------------------------------------------------

class TestBuildToFile:
    def test_build_to_file(self, minimal_schema, tmp_path):
        builder = PPTXBuilder(minimal_schema)
        output = tmp_path / "test.pptx"
        builder.build_to_file({}, output)
        assert output.exists()
        assert output.stat().st_size > 0
        prs = Presentation(str(output))
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Full 14-slide integration test
# ---------------------------------------------------------------------------

class TestFullPipeline:
    def _sample_payload(self):
        """Build a representative payload with data for all 14 slides."""
        return {
            # Cover KPIs
            "cover.report_title": "No7 US Monthly eComm Report",
            "cover.report_period": "January 2026 Overview",
            "cover.total_revenue": 1234567,
            "cover.total_orders": 12345,
            "cover.aov": 100.0,
            "cover.new_customers": 4500,
            "cover.cvr": 3.6,
            "cover.cos": 12.5,
            "cover.revenue_vs_target": 5.2,
            "cover.orders_vs_target": 3.1,
            "cover.aov_vs_target": -1.2,
            "cover.nc_vs_target": 8.0,
            "cover.cvr_vs_target": 0.5,
            "cover.cos_vs_target": -0.3,
            # TOC
            "toc.items": [
                "eComm Performance Overview",
                "Daily Performance",
                "Promotion Performance",
                "Product Performance",
                "Channel Deep Dives",
                "Outlook",
            ],
            # Dividers
            "divider.ecomm_title": "eComm Performance",
            "divider.channels_title": "Channel Deep Dives",
            "divider.outlook_title": "Outlook",
            # Executive summary
            "exec.title": "Executive Summary",
            "exec.performance_rows": [
                {
                    "channel": "Total",
                    "revenue": 1234567,
                    "revenue_vs_target": 5.2,
                    "revenue_vs_ly": 12.3,
                    "orders": 12345,
                    "sessions": 345678,
                    "cvr": 3.6,
                    "aov": 100.0,
                    "cos": 12.5,
                    "new_customers": 4500,
                },
                {
                    "channel": "DIRECT",
                    "revenue": 400000,
                    "revenue_vs_target": 8.1,
                    "revenue_vs_ly": 15.0,
                    "orders": 4000,
                    "sessions": 120000,
                    "cvr": 3.3,
                    "aov": 100.0,
                    "cos": 10.0,
                    "new_customers": 1500,
                },
            ],
            "exec.narrative": "Strong month driven by DIRECT channel outperformance.",
            # Daily performance
            "daily.title": "Daily Performance",
            "daily.dates": ["1/1", "1/2", "1/3", "1/4", "1/5"],
            "daily.revenue_actual": [40000, 45000, 38000, 52000, 48000],
            "daily.revenue_target": [42000, 42000, 42000, 42000, 42000],
            "daily.revenue_ly": [35000, 38000, 32000, 45000, 41000],
            "daily.campaign_rows": [
                {"date": "1/1", "activity": "New Year Sale Launch"},
                {"date": "1/3", "activity": "Email Blast - Winter"},
            ],
            "daily.revenue_achieved_pct": 75.0,
            "daily.revenue_remaining_pct": 25.0,
            # Promotions
            "promo.title": "Promotion Performance",
            "promo.rows": [
                {
                    "promotion_name": "New Year Sale",
                    "channel": "All",
                    "redemptions": 5000,
                    "redemptions_vs_ly": 12.5,
                    "revenue": 250000,
                    "revenue_vs_ly": 8.3,
                    "discount_amount": 45000,
                },
            ],
            # Products
            "product.title": "Product Performance",
            "product.rows": [
                {
                    "product_name": "No7 Serum",
                    "units": 3500,
                    "units_vs_ly": 15.2,
                    "revenue": 175000,
                    "revenue_vs_ly": 18.1,
                    "aov": 50.0,
                    "avg_selling_price": 50.0,
                    "discount_pct": 5.0,
                    "new_customers": 800,
                },
            ],
            # CRM
            "crm.title": "CRM Performance",
            "crm.emails_sent": 250000,
            "crm.emails_sent_vs_ly": 5.0,
            "crm.open_rate": 22.5,
            "crm.open_rate_vs_ly": 1.2,
            "crm.ctr": 3.8,
            "crm.ctr_vs_ly": -0.5,
            "crm.revenue": 180000,
            "crm.revenue_vs_ly": 12.0,
            "crm.cvr": 4.2,
            "crm.cvr_vs_ly": 0.3,
            "crm.aov": 95.0,
            "crm.aov_vs_ly": -2.1,
            "crm.detail_rows": [
                {
                    "campaign_type": "Manual Campaigns",
                    "emails_sent": 150000,
                    "open_rate": 25.0,
                    "ctr": 4.2,
                    "sessions": 6300,
                    "orders": 252,
                    "cvr": 4.0,
                    "revenue": 120000,
                    "aov": 476.19,
                    "revenue_vs_ly": 10.0,
                },
            ],
            # Affiliate
            "affiliate.title": "Affiliate Performance",
            "affiliate.revenue": 95000,
            "affiliate.revenue_vs_ly": 7.5,
            "affiliate.cos": 8.0,
            "affiliate.cos_vs_ly": -1.0,
            "affiliate.roas": 12.5,
            "affiliate.roas_vs_ly": 1.2,
            "affiliate.orders": 950,
            "affiliate.orders_vs_ly": 5.0,
            "affiliate.cvr": 2.8,
            "affiliate.cvr_vs_ly": 0.2,
            "affiliate.publisher_rows": [
                {
                    "publisher_name": "Publisher A",
                    "revenue": 30000,
                    "revenue_vs_ly": 10.0,
                    "commission": 2400,
                    "cos": 8.0,
                    "orders": 300,
                    "cvr": 3.0,
                    "sessions": 10000,
                    "aov": 100.0,
                },
            ],
            # SEO
            "seo.title": "SEO Performance",
            "seo.revenue": 120000,
            "seo.revenue_vs_ly": 6.0,
            "seo.sessions": 80000,
            "seo.sessions_vs_ly": 4.5,
            "seo.cvr": 3.0,
            "seo.cvr_vs_ly": 0.1,
            "seo.orders": 2400,
            "seo.orders_vs_ly": 6.5,
            "seo.aov": 50.0,
            "seo.aov_vs_ly": -0.5,
            "seo.narrative": "Organic traffic grew steadily. Focus on content strategy.",
            # Upcoming promotions
            "upcoming.title": "Upcoming Promotions",
            "upcoming.rows": [
                {
                    "date": "Feb 1-14",
                    "promotion": "Valentine's Day Sale",
                    "discount": "20% off",
                    "channels": "All Channels",
                },
            ],
            # Next steps
            "next_steps.title": "Next Steps",
            "next_steps.items": "1. Review Feb targets\n2. Launch Valentine campaign",
        }

    def test_full_14_slide_build(self, full_schema):
        payload = self._sample_payload()
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)

        assert len(prs.slides) == 14

        # Verify slide dimensions
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_full_build_empty_payload(self, full_schema):
        """Building with empty payload should not crash."""
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build({})
        prs = _bytes_to_prs(pptx_bytes)
        assert len(prs.slides) == 14

    def test_full_build_cover_content(self, full_schema):
        payload = self._sample_payload()
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)

        # Cover slide (index 0) should contain formatted KPIs
        cover = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in cover.shapes if s.has_text_frame
        )
        assert "$1.2m" in all_text  # Revenue
        assert "Revenue" in all_text
        assert "No7 US Monthly eComm Report" in all_text

    def test_full_build_table_slide(self, full_schema):
        payload = self._sample_payload()
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)

        # Executive summary slide (index 3) should have a table
        exec_slide = prs.slides[3]
        tables = [s for s in exec_slide.shapes if s.has_table]
        assert len(tables) >= 1

    def test_full_build_chart_slide(self, full_schema):
        payload = self._sample_payload()
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)

        # Daily performance slide (index 4) should have charts
        daily_slide = prs.slides[4]
        charts = [s for s in daily_slide.shapes if s.has_chart]
        assert len(charts) >= 1

    def test_full_build_divider_slides(self, full_schema):
        payload = self._sample_payload()
        builder = PPTXBuilder(full_schema)
        pptx_bytes = builder.build(payload)
        prs = _bytes_to_prs(pptx_bytes)

        # Divider slides (indices 2, 7, 11) should have background fill
        for idx in [2, 7, 11]:
            slide = prs.slides[idx]
            bg = slide.background
            assert bg.fill.fore_color.rgb == RGBColor(0x00, 0x65, 0xE0)

    def test_full_build_to_file(self, full_schema, tmp_path):
        payload = self._sample_payload()
        builder = PPTXBuilder(full_schema)
        output = tmp_path / "full_report.pptx"
        builder.build_to_file(payload, output)
        assert output.exists()
        prs = Presentation(str(output))
        assert len(prs.slides) == 14
